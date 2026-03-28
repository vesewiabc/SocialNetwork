from flask import Flask, render_template, request, redirect, session, flash, url_for, jsonify
import sqlite3
import os
from datetime import datetime, timedelta
import hashlib
from werkzeug.utils import secure_filename
import re
import requests
import random
import string
import threading
import asyncio
import logging

# ======================== НАСТРОЙКИ TELEGRAM БОТА ========================
BOT_TOKEN = "8542566873:AAGCjyU1Q5IM4tip_MC77Jt43lHh8eK7Bbk"  # <-- вставь токен от @BotFather
BOT_USERNAME = "LinkA_2FA_Bot"  # <-- username бота без @, например: my2fa_bot
# =========================================================================

# Импорты для просмотра документов
try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl

    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

app = Flask(__name__)
app.secret_key = "123"

# Настройки для загрузки файлов
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'static', 'uploads', 'avatars')
GROUP_UPLOAD_FOLDER = os.path.join(BASE_DIR, 'static', 'uploads', 'groups')
POST_MEDIA_FOLDER = os.path.join(BASE_DIR, 'static', 'uploads', 'posts')
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
ALLOWED_VIDEO_EXTENSIONS = {'mp4', 'avi', 'mov', 'wmv', 'mkv', 'webm'}
ALLOWED_DOCUMENT_EXTENSIONS = {'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx', 'txt', 'zip', 'rar', '7z'}
ALLOWED_EXTENSIONS = ALLOWED_IMAGE_EXTENSIONS | ALLOWED_VIDEO_EXTENSIONS | ALLOWED_DOCUMENT_EXTENSIONS

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['GROUP_UPLOAD_FOLDER'] = GROUP_UPLOAD_FOLDER
app.config['POST_MEDIA_FOLDER'] = POST_MEDIA_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB

# Создаем папки для загрузок, если их нет
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GROUP_UPLOAD_FOLDER, exist_ok=True)
os.makedirs(POST_MEDIA_FOLDER, exist_ok=True)


def get_current_date():
    """Возвращает текущую дату в формате ДД.ММ.ГГГГ"""
    return datetime.now().strftime('%d.%m.%Y')


def get_current_datetime():
    """Возвращает текущую дату и время в формате ДД.ММ.ГГГГ ЧЧ:ММ"""
    return datetime.now().strftime('%d.%m.%Y %H:%M')


def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def allowed_image_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_IMAGE_EXTENSIONS


def allowed_video_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_VIDEO_EXTENSIONS


# Добавляем соединение с БД для работы со строками как словарями
def get_db_connection():
    conn = sqlite3.connect('users.db')
    conn.row_factory = sqlite3.Row
    return conn


# Функция для преобразования Row в словарь
def row_to_dict(row):
    if row is None:
        return None
    return dict(row)


# Функция для преобразования списка Row в список словарей
def rows_to_dicts(rows):
    return [dict(row) for row in rows]


# ==================== 2FA ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ====================

def generate_code(length=6) -> str:
    """Генерирует случайный 6-значный цифровой код для входа"""
    return ''.join(random.choices(string.digits, k=length))


def generate_link_code(length=32) -> str:
    """Генерирует случайный код для привязки Telegram"""
    return ''.join(random.choices(string.ascii_letters + string.digits, k=length))


def send_telegram_2fa_code(chat_id: str, code: str, username: str) -> bool:
    """Отправляет код 2FA через Telegram Bot API (синхронный HTTP-запрос)"""
    url = f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage"
    text = (
        f"🔐 *Код для входа на сайт*\n\n"
        f"Аккаунт: *{username}*\n"
        f"Код: `{code}`\n\n"
        f"⏱ Код действует *5 минут*.\n"
        f"Если вы не пытались войти — игнорируйте это сообщение."
    )
    try:
        resp = requests.post(url, json={
            "chat_id": chat_id,
            "text": text,
            "parse_mode": "Markdown"
        }, timeout=10)
        return resp.status_code == 200
    except Exception as e:
        print(f"[2FA] Ошибка отправки Telegram: {e}")
        return False


# Функция для добавления недостающих столбцов в таблицы
def migrate_database():
    conn = get_db_connection()
    cursor = conn.cursor()

    # Проверяем и добавляем недостающие столбцы
    try:
        cursor.execute("SELECT role FROM users LIMIT 1")
    except sqlite3.OperationalError:
        print("Добавляем столбец role в таблицу users...")
        cursor.execute('ALTER TABLE users ADD COLUMN role TEXT DEFAULT "user"')

    try:
        cursor.execute("SELECT is_banned FROM users LIMIT 1")
    except sqlite3.OperationalError:
        print("Добавляем столбец is_banned в таблицу users...")
        cursor.execute('ALTER TABLE users ADD COLUMN is_banned INTEGER DEFAULT 0')

    try:
        cursor.execute("SELECT avatar FROM user_profiles LIMIT 1")
    except sqlite3.OperationalError:
        print("Добавляем столбец avatar в таблицу user_profiles...")
        cursor.execute('ALTER TABLE user_profiles ADD COLUMN avatar TEXT')
        cursor.execute('UPDATE user_profiles SET avatar = "default_avatar.png" WHERE avatar IS NULL')

    # Проверяем и добавляем request_permissions если нет
    try:
        cursor.execute("SELECT request_permissions FROM groups LIMIT 1")
    except sqlite3.OperationalError:
        print("Добавляем столбец request_permissions в таблицу groups...")
        cursor.execute('ALTER TABLE groups ADD COLUMN request_permissions TEXT DEFAULT "moderators"')

    # Проверяем и добавляем таблицу post_media если нет
    try:
        cursor.execute("SELECT post_id FROM post_media LIMIT 1")
    except sqlite3.OperationalError:
        print("Таблица post_media не существует, создаем...")
        cursor.execute('''
                CREATE TABLE post_media (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    post_id INTEGER,
                    group_post_id INTEGER,
                    filename TEXT NOT NULL,
                    file_type TEXT NOT NULL,
                    thumbnail TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    FOREIGN KEY (post_id) REFERENCES posts(id) ON DELETE CASCADE,
                    FOREIGN KEY (group_post_id) REFERENCES group_posts(id) ON DELETE CASCADE
                )
            ''')
    else:
        # Проверяем структуру
        cursor.execute("PRAGMA table_info(post_media)")
        columns = cursor.fetchall()
        for col in columns:
            if col[1] == 'post_id' and col[3] == 1:  # post_id имеет NOT NULL
                print("Исправляем структуру таблицы post_media...")
                # Создаем временную таблицу
                cursor.execute('''
                        CREATE TABLE post_media_new (
                            id INTEGER PRIMARY KEY AUTOINCREMENT,
                            post_id INTEGER,
                            group_post_id INTEGER,
                            filename TEXT NOT NULL,
                            file_type TEXT NOT NULL,
                            thumbnail TEXT,
                            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                            FOREIGN KEY (post_id) REFERENCES posts(id) ON DELETE CASCADE,
                            FOREIGN KEY (group_post_id) REFERENCES group_posts(id) ON DELETE CASCADE
                        )
                    ''')

                # Копируем данные
                cursor.execute('''
                        INSERT INTO post_media_new 
                        SELECT id, post_id, group_post_id, filename, file_type, thumbnail, created_at 
                        FROM post_media
                    ''')

                # Удаляем старую и переименовываем новую
                cursor.execute('DROP TABLE post_media')
                cursor.execute('ALTER TABLE post_media_new RENAME TO post_media')
                break

    # Добавляем колонку original_filename если её нет
    cursor.execute("PRAGMA table_info(post_media)")
    pm_columns = [col[1] for col in cursor.fetchall()]
    if 'original_filename' not in pm_columns:
        cursor.execute('ALTER TABLE post_media ADD COLUMN original_filename TEXT')
        print("Добавлена колонка original_filename в post_media")

    # Создаем таблицу черного списка, если ее нет
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS blacklist (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        blocker_id INTEGER NOT NULL,
        blocked_id INTEGER NOT NULL,
        reason TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (blocker_id) REFERENCES users(id),
        FOREIGN KEY (blocked_id) REFERENCES users(id),
        UNIQUE(blocker_id, blocked_id)
    )
    ''')

    # Создаем таблицу двухфакторной аутентификации
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS two_factor_auth (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL UNIQUE,
        telegram_chat_id TEXT,
        is_enabled INTEGER DEFAULT 0,
        link_code TEXT,
        link_code_expires TIMESTAMP,
        auth_code TEXT,
        auth_code_expires TIMESTAMP,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
    )
    ''')

    conn.commit()
    conn.close()


# Таблицы для хранения пользователей и их данных
def ensure_site_news_table():
    conn = get_db_connection()
    conn.execute('''
        CREATE TABLE IF NOT EXISTS site_news (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            body TEXT NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()



def create_tables():
    connection = get_db_connection()
    cursor = connection.cursor()

    # Таблица пользователей
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT NOT NULL UNIQUE,
        password TEXT NOT NULL,
        role TEXT DEFAULT 'user',
        is_banned INTEGER DEFAULT 0,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    ''')

    # Таблица профилей пользователей
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS user_profiles (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL UNIQUE,
        full_name TEXT,
        bio TEXT,
        location TEXT,
        website TEXT,
        avatar TEXT DEFAULT 'default_avatar.png',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id)
    )
    ''')

    # Таблица личных постов
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS posts (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        content TEXT NOT NULL,
        visibility TEXT DEFAULT 'public',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id)
    )
    ''')

    # Таблица друзей
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS friendships (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        sender_id INTEGER NOT NULL,
        receiver_id INTEGER NOT NULL,
        status TEXT DEFAULT 'pending',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (sender_id) REFERENCES users(id),
        FOREIGN KEY (receiver_id) REFERENCES users(id),
        UNIQUE(sender_id, receiver_id)
    )
    ''')

    # Таблица лайков постов
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS post_likes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        post_id INTEGER NOT NULL,
        user_id INTEGER NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (post_id) REFERENCES posts(id),
        FOREIGN KEY (user_id) REFERENCES users(id),
        UNIQUE(post_id, user_id)
    )
    ''')

    # Таблица комментариев
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS comments (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        post_id INTEGER NOT NULL,
        user_id INTEGER NOT NULL,
        content TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (post_id) REFERENCES posts(id),
        FOREIGN KEY (user_id) REFERENCES users(id)
    )
    ''')

    # Таблица групп/пабликов
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS groups (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT NOT NULL,
        description TEXT,
        avatar TEXT DEFAULT 'default_group.png',
        creator_id INTEGER NOT NULL,
        is_public BOOLEAN DEFAULT 1,
        post_permissions TEXT DEFAULT 'all', -- 'admins', 'moderators', 'all'
        request_permissions TEXT DEFAULT 'moderators', -- 'admins', 'moderators'
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (creator_id) REFERENCES users(id)
    )
    ''')

    # Таблица подписчиков групп
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS group_members (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        group_id INTEGER NOT NULL,
        user_id INTEGER NOT NULL,
        role TEXT DEFAULT 'member', -- 'admin', 'moderator', 'member'
        joined_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (group_id) REFERENCES groups(id),
        FOREIGN KEY (user_id) REFERENCES users(id),
        UNIQUE(group_id, user_id)
    )
    ''')

    # Таблица постов в группах
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS group_posts (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        group_id INTEGER NOT NULL,
        author_id INTEGER NOT NULL,
        content TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (group_id) REFERENCES groups(id),
        FOREIGN KEY (author_id) REFERENCES users(id)
    )
    ''')

    # Таблица лайков постов в группах
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS group_post_likes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        post_id INTEGER NOT NULL,
        user_id INTEGER NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (post_id) REFERENCES group_posts(id),
        FOREIGN KEY (user_id) REFERENCES users(id),
        UNIQUE(post_id, user_id)
    )
    ''')

    # Таблица заявок на вступление в приватные группы
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS group_requests (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        group_id INTEGER NOT NULL,
        user_id INTEGER NOT NULL,
        status TEXT DEFAULT 'pending', -- 'pending', 'approved', 'rejected'
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (group_id) REFERENCES groups(id),
        FOREIGN KEY (user_id) REFERENCES users(id),
        UNIQUE(group_id, user_id)
    )
    ''')

    # Таблица медиафайлов для постов
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS post_media (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        post_id INTEGER,
        group_post_id INTEGER,
        filename TEXT NOT NULL,
        file_type TEXT NOT NULL, -- 'image' или 'video'
        thumbnail TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (post_id) REFERENCES posts(id) ON DELETE CASCADE,
        FOREIGN KEY (group_post_id) REFERENCES group_posts(id) ON DELETE CASCADE
    )
    ''')

    # Таблица импортированных новостей
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS imported_news (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT NOT NULL,
        description TEXT,
        link TEXT NOT NULL UNIQUE,
        source TEXT DEFAULT 'RBC',
        published TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        imported_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    ''')

    # Таблица жалоб
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS reports (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        reporter_id INTEGER NOT NULL,
        reported_id INTEGER NOT NULL,
        reason TEXT NOT NULL,
        status TEXT DEFAULT 'pending',
        admin_notes TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (reporter_id) REFERENCES users(id),
        FOREIGN KEY (reported_id) REFERENCES users(id)
    )
    ''')

    connection.commit()
    connection.close()

    # Создаём таблицу новостей соцсети
    ensure_site_news_table()

    # Выполняем миграции для существующих таблиц
    migrate_database()

    # Создаем тех-админа, если его нет
    create_tech_admin()


def create_tech_admin():
    conn = get_db_connection()
    cursor = conn.cursor()

    # Проверяем, существует ли уже тех-админ
    cursor.execute("SELECT id FROM users WHERE username = 'techadmin'")
    tech_admin = cursor.fetchone()

    if not tech_admin:
        # Создаем тех-админа
        cursor.execute("INSERT INTO users (username, password, role) VALUES (?, ?, ?)",
                       ('techadmin', 'techadmin123', 'techadmin'))
        print("Создан аккаунт тех-админа: Логин: techadmin, Пароль: techadmin123")

    # Проверяем, существует ли админ
    cursor.execute("SELECT id FROM users WHERE username = 'admin'")
    admin = cursor.fetchone()

    if not admin:
        # Обновляем существующего админа или создаем нового
        cursor.execute("SELECT id FROM users WHERE username = 'admin'")
        existing_admin = cursor.fetchone()

        if existing_admin:
            # Обновляем роль существующего админа
            cursor.execute("UPDATE users SET role = 'admin' WHERE username = 'admin'")
        else:
            # Создаем нового админа
            cursor.execute("INSERT INTO users (username, password, role) VALUES (?, ?, ?)",
                           ('admin', 'admin123', 'admin'))
            print("Создан аккаунт администратора: Логин: admin, Пароль: admin123")

    conn.commit()
    conn.close()


# Добавляем функцию для форматирования даты в шаблонах
@app.context_processor
def utility_processor():
    def format_datetime(value, format='%d.%m.%Y %H:%M'):
        if not value:
            return ''

        # Если дата уже в формате ДД.ММ.ГГГГ, просто добавляем время если нужно
        if isinstance(value, str):
            if len(value) == 10 and '.' in value:  # ДД.ММ.ГГГГ
                if format == '%d.%m.%Y':
                    return value
                else:
                    return f"{value} 00:00"  # Добавляем время по умолчанию

        return value

    def get_default_avatar(avatar_type='user'):
        """Возвращает путь к дефолтной аватарке"""
        if avatar_type == 'group':
            return url_for('static', filename='defaults/for_groups.png')
        else:  # user
            return url_for('static', filename='defaults/for_users.png')

    def get_current_user_avatar():
        """Возвращает аватарку текущего пользователя для шапки"""
        user_id = session.get('user_id')
        if not user_id:
            return None
        try:
            conn = sqlite3.connect('users.db')
            cursor = conn.cursor()
            cursor.execute("SELECT avatar FROM user_profiles WHERE user_id = ?", (user_id,))
            row = cursor.fetchone()
            conn.close()
            if row and row[0]:
                return row[0]
        except:
            pass
        return None

    def get_current_user_info():
        """Возвращает имя и юзернейм текущего пользователя для шапки"""
        user_id = session.get('user_id')
        if not user_id:
            return None
        try:
            conn = sqlite3.connect('users.db')
            cursor = conn.cursor()
            cursor.execute("""
                SELECT u.username, p.full_name
                FROM users u
                LEFT JOIN user_profiles p ON p.user_id = u.id
                WHERE u.id = ?
            """, (user_id,))
            row = cursor.fetchone()
            conn.close()
            if row:
                return {'username': row[0], 'full_name': row[1] or row[0]}
        except:
            pass
        return None

    return {
        'datetime_format': format_datetime,
        'default_avatar': get_default_avatar,
        'current_user_avatar': get_current_user_avatar,
        'current_user_info': get_current_user_info
    }


# Функция для получения медиафайлов поста
def get_post_media(post_id, is_group_post=True):
    """Получает медиафайлы для поста"""
    conn = get_db_connection()

    if is_group_post:
        media = conn.execute('''
            SELECT id, filename, file_type, thumbnail, original_filename
            FROM post_media 
            WHERE group_post_id = ?
            ORDER BY id
        ''', (post_id,)).fetchall()
    else:
        media = conn.execute('''
            SELECT id, filename, file_type, thumbnail, original_filename
            FROM post_media 
            WHERE post_id = ?
            ORDER BY id
        ''', (post_id,)).fetchall()

    conn.close()
    return rows_to_dicts(media)


# ==================== ФУНКЦИИ ДЛЯ КОММЕНТАРИЕВ И ЛАЙКОВ ====================

def get_post_comments(post_id):
    """Получение комментариев для поста"""
    conn = get_db_connection()
    comments = rows_to_dicts(conn.execute('''
        SELECT c.*, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar
        FROM comments c
        JOIN users u ON c.user_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE c.post_id = ?
        ORDER BY c.created_at DESC
    ''', (post_id,)).fetchall())
    conn.close()
    return comments


def get_post_likes(post_id):
    """Получение лайков для поста"""
    conn = get_db_connection()
    likes = rows_to_dicts(conn.execute('''
        SELECT pl.*, u.username
        FROM post_likes pl
        JOIN users u ON pl.user_id = u.id
        WHERE pl.post_id = ?
    ''', (post_id,)).fetchall())
    conn.close()
    return likes


def has_user_liked_post(post_id, user_id):
    """Проверка, поставил ли пользователь лайк посту"""
    conn = get_db_connection()
    result = conn.execute('''
        SELECT id FROM post_likes 
        WHERE post_id = ? AND user_id = ?
    ''', (post_id, user_id)).fetchone()
    conn.close()
    return result is not None


# ==================== НОВЫЕ МАРШРУТЫ ====================

@app.route('/edit_post/<int:post_id>', methods=['GET', 'POST'])
def edit_post(post_id):
    """Редактирование поста"""
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем пост
    post = conn.execute('SELECT * FROM posts WHERE id = ?', (post_id,)).fetchone()

    if not post:
        conn.close()
        flash('Пост не найден', 'error')
        return redirect('/my_posts')

    # Проверяем права доступа
    if post['user_id'] != user_id:
        conn.close()
        flash('Вы не можете редактировать этот пост', 'error')
        return redirect('/my_posts')

    if request.method == 'POST':
        content = request.form.get('content', '').strip()

        if not content:
            flash('Пост не может быть пустым', 'error')
            return redirect(f'/edit_post/{post_id}')

        # Обновляем пост
        conn.execute('UPDATE posts SET content = ? WHERE id = ?', (content, post_id))
        conn.commit()
        conn.close()

        flash('Пост успешно обновлен!', 'success')
        return redirect('/my_posts')

    # GET запрос - показываем форму редактирования
    conn.close()
    return render_template('edit_post.html', post=dict(post))


@app.route('/faq')
def faq():
    """Страница часто задаваемых вопросов"""
    if 'user_id' not in session:
        return redirect('/login')

    username = session.get('username', 'Пользователь')
    return render_template('faq.html', username=username)


@app.route('/delete_post_route/<int:post_id>', methods=['POST'])
def delete_post_route(post_id):
    """Удаление поста (POST запрос)"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем пост
    post = conn.execute('SELECT * FROM posts WHERE id = ?', (post_id,)).fetchone()

    if not post:
        conn.close()
        return jsonify({'success': False, 'error': 'Пост не найден'})

    # Проверяем права доступа
    if post['user_id'] != user_id:
        conn.close()
        return jsonify({'success': False, 'error': 'Нет прав на удаление'})

    try:
        # Удаляем лайки и комментарии
        conn.execute('DELETE FROM post_likes WHERE post_id = ?', (post_id,))
        conn.execute('DELETE FROM comments WHERE post_id = ?', (post_id,))

        # Удаляем медиафайлы
        media_files = conn.execute('SELECT filename FROM post_media WHERE post_id = ?', (post_id,)).fetchall()
        for media in media_files:
            try:
                os.remove(os.path.join(app.config['POST_MEDIA_FOLDER'], media['filename']))
            except:
                pass

        # Удаляем записи о медиа
        conn.execute('DELETE FROM post_media WHERE post_id = ?', (post_id,))

        # Удаляем сам пост
        conn.execute('DELETE FROM posts WHERE id = ?', (post_id,))
        conn.commit()
        conn.close()

        return jsonify({'success': True, 'message': 'Пост удален'})

    except Exception as e:
        conn.rollback()
        conn.close()
        return jsonify({'success': False, 'error': str(e)})


@app.route('/like_post_action/<int:post_id>', methods=['POST'])
def like_post_action(post_id):
    """Поставить/убрать лайк посту"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    conn = get_db_connection()

    try:
        # Проверяем, существует ли пост
        post = conn.execute('SELECT id FROM posts WHERE id = ?', (post_id,)).fetchone()
        if not post:
            conn.close()
            return jsonify({'success': False, 'error': 'Пост не найден'})

        # Проверяем, лайкал ли уже пользователь
        existing_like = conn.execute('''
            SELECT id FROM post_likes WHERE post_id = ? AND user_id = ?
        ''', (post_id, user_id)).fetchone()

        if existing_like:
            # Убираем лайк
            conn.execute('DELETE FROM post_likes WHERE id = ?', (existing_like['id'],))
            action = 'unliked'
        else:
            # Ставим лайк
            conn.execute('INSERT INTO post_likes (post_id, user_id) VALUES (?, ?)', (post_id, user_id))
            action = 'liked'

        # Получаем новое количество лайков
        likes_count = conn.execute('SELECT COUNT(*) as count FROM post_likes WHERE post_id = ?', (post_id,)).fetchone()[
            'count']

        conn.commit()
        conn.close()

        return jsonify({
            'success': True,
            'action': action,
            'likes_count': likes_count,
            'liked': action == 'liked'
        })

    except Exception as e:
        conn.rollback()
        conn.close()
        return jsonify({'success': False, 'error': str(e)})


@app.route('/add_comment/<int:post_id>', methods=['POST'])
def add_comment(post_id):
    """Добавить комментарий к посту"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    data = request.get_json()
    content = data.get('content', '').strip()

    if not content:
        return jsonify({'success': False, 'error': 'Комментарий не может быть пустым'})

    conn = get_db_connection()

    try:
        # Проверяем, существует ли пост
        post = conn.execute('SELECT id FROM posts WHERE id = ?', (post_id,)).fetchone()
        if not post:
            conn.close()
            return jsonify({'success': False, 'error': 'Пост не найден'})

        # Добавляем комментарий - ИСПОЛЬЗУЕМ CURSOR
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO comments (post_id, user_id, content)
            VALUES (?, ?, ?)
        ''', (post_id, user_id, content))

        # Получаем ID нового комментария через cursor.lastrowid
        comment_id = cursor.lastrowid

        # Получаем информацию о комментаторе для ответа
        comment_data = conn.execute('''
            SELECT c.*, u.username, up.full_name, 
                   COALESCE(up.avatar, 'default_avatar.png') as avatar
            FROM comments c
            JOIN users u ON c.user_id = u.id
            LEFT JOIN user_profiles up ON u.id = up.user_id
            WHERE c.id = ?
        ''', (comment_id,)).fetchone()

        # Получаем общее количество комментариев
        comments_count = \
            conn.execute('SELECT COUNT(*) as count FROM comments WHERE post_id = ?', (post_id,)).fetchone()['count']

        conn.commit()
        conn.close()

        return jsonify({
            'success': True,
            'comment': dict(comment_data),
            'comments_count': comments_count
        })

    except Exception as e:
        conn.rollback()
        conn.close()
        return jsonify({'success': False, 'error': str(e)})


@app.route('/delete_comment/<int:comment_id>', methods=['POST'])
def delete_comment(comment_id):
    """Удалить комментарий"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    conn = get_db_connection()

    try:
        # Проверяем, существует ли комментарий
        comment = conn.execute('SELECT * FROM comments WHERE id = ?', (comment_id,)).fetchone()
        if not comment:
            conn.close()
            return jsonify({'success': False, 'error': 'Комментарий не найден'})

        # Проверяем права доступа (только автор комментария или автор поста)
        if comment['user_id'] != user_id:
            # Проверяем, является ли пользователь автором поста
            post = conn.execute('SELECT user_id FROM posts WHERE id = ?', (comment['post_id'],)).fetchone()
            if not post or post['user_id'] != user_id:
                conn.close()
                return jsonify({'success': False, 'error': 'Нет прав на удаление'})

        # Получаем post_id перед удалением
        post_id = comment['post_id']

        # Удаляем комментарий
        conn.execute('DELETE FROM comments WHERE id = ?', (comment_id,))

        # Получаем новое количество комментариев
        comments_count = \
            conn.execute('SELECT COUNT(*) as count FROM comments WHERE post_id = ?', (post_id,)).fetchone()['count']

        conn.commit()
        conn.close()

        return jsonify({
            'success': True,
            'comments_count': comments_count
        })

    except Exception as e:
        conn.rollback()
        conn.close()
        return jsonify({'success': False, 'error': str(e)})


@app.route('/post/<int:post_id>')
def view_post(post_id):
    """Просмотр отдельного поста с комментариями"""
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем пост
    post = conn.execute('''
        SELECT p.*, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar
        FROM posts p
        JOIN users u ON p.user_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE p.id = ?
    ''', (post_id,)).fetchone()

    if not post:
        conn.close()
        flash('Пост не найден', 'error')
        return redirect('/home')

    # Получаем комментарии
    comments = get_post_comments(post_id)

    # Получаем лайки
    likes = get_post_likes(post_id)
    has_liked = has_user_liked_post(post_id, user_id)

    # Получаем количество лайков и комментариев
    likes_count = len(likes)
    comments_count = len(comments)

    conn.close()

    return render_template('view_post.html',
                           post=dict(post),
                           comments=comments,
                           likes=likes,
                           likes_count=likes_count,
                           comments_count=comments_count,
                           has_liked=has_liked,
                           current_user_id=user_id)


@app.template_filter('format_date')
def format_date_filter(value):
    if not value:
        return ''

    if isinstance(value, str):
        try:
            # Пробуем разные форматы входящих дат
            formats_to_try = [
                '%Y-%m-%d',  # 2026-01-27
                '%Y-%m-%d %H:%M:%S',  # 2026-01-27 10:30:15
                '%Y-%m-%d %H:%M:%S.%f',  # 2026-01-27 10:30:15.123456
                '%Y.%m.%d',  # 2026.01.27 - ВАШ ФОРМАТ
                '%Y.%m.%d %H:%M:%S',  # 2026.01.27 10:30:15
                '%d-%m-%Y %H:%M:%S',  # 27-01-2026 10:30:15
                '%d-%m-%Y %H:%M:%S.%f',  # 27-01-2026 10:30:15.123456
            ]

            for fmt in formats_to_try:
                try:
                    dt = datetime.strptime(value, fmt)
                    return dt.strftime('%d.%m.%Y')
                except:
                    continue
        except Exception as e:
            print(f"Ошибка форматирования даты {value}: {e}")

    # Если значение - объект datetime
    if hasattr(value, 'strftime'):
        return value.strftime('%d.%m.%Y')

    # Если ничего не сработало, возвращаем как есть
    return value


# Проверка прав на создание поста
def check_post_permission(group_id, user_id, conn):
    """Проверяет, может ли пользователь создавать посты в группе"""
    # Получаем настройки группы
    group = conn.execute('SELECT post_permissions FROM groups WHERE id = ?', (group_id,)).fetchone()

    if not group:
        return False

    # Получаем роль пользователя
    membership = conn.execute('SELECT role FROM group_members WHERE group_id = ? AND user_id = ?',
                              (group_id, user_id)).fetchone()

    if not membership:
        return False

    user_role = membership['role']
    post_permissions = group['post_permissions']

    # Проверяем права в зависимости от настроек
    if post_permissions == 'admins':
        return user_role == 'admin'
    elif post_permissions == 'moderators':
        return user_role in ['admin', 'moderator']
    else:  # 'all'
        return True


# ==================== НОВОСТИ СОЦСЕТИ ====================

def get_posts_feed(user_id, limit=20, filter_type='all', offset=0):
    """Получение ленты постов пользователей (без импортированных новостей)"""
    conn = get_db_connection()
    try:
        if filter_type == 'mine':
            filter_clause = 'WHERE p.user_id = :uid'
            params = {'uid': user_id, 'limit': limit, 'offset': offset}
        elif filter_type == 'friends':
            filter_clause = '''WHERE p.user_id IN (
                SELECT CASE WHEN sender_id = :uid THEN receiver_id ELSE sender_id END
                FROM friendships WHERE (sender_id = :uid OR receiver_id = :uid) AND status = 'accepted'
            )'''
            params = {'uid': user_id, 'limit': limit, 'offset': offset}
        else:
            filter_clause = ''
            params = {'limit': limit, 'offset': offset}

        rows = conn.execute(f'''
            SELECT p.*, u.username,
                   COALESCE(up.full_name, u.username) as author_name,
                   COALESCE(up.avatar, '') as author_avatar,
                   (SELECT COUNT(*) FROM post_likes WHERE post_id = p.id) as likes_count,
                   (SELECT COUNT(*) FROM comments WHERE post_id = p.id) as comments_count,
                   EXISTS(SELECT 1 FROM post_likes WHERE post_id = p.id AND user_id = :self_uid) as is_liked
            FROM posts p
            JOIN users u ON p.user_id = u.id
            LEFT JOIN user_profiles up ON up.user_id = p.user_id
            {filter_clause}
            ORDER BY p.created_at DESC
            LIMIT :limit OFFSET :offset
        ''', {**params, 'self_uid': user_id}).fetchall()

        posts = rows_to_dicts(rows)
        for post in posts:
            post['type'] = 'post'
            media_files = conn.execute(
                'SELECT filename, file_type, original_filename FROM post_media WHERE post_id = ? ORDER BY id',
                (post['id'],)
            ).fetchall()
            post['media_files'] = rows_to_dicts(media_files) if media_files else []
    except Exception as e:
        print(f"Ошибка get_posts_feed: {e}")
        posts = []
    finally:
        conn.close()
    return posts


@app.route('/admin/news/add', methods=['POST'])
def admin_add_news():
    if 'user_id' not in session or session.get('role') != 'admin':
        return jsonify({'success': False, 'error': 'Нет доступа'}), 403
    ensure_site_news_table()
    data = request.get_json()
    title = (data.get('title') or '').strip()
    body = (data.get('body') or '').strip()
    if not title or not body:
        return jsonify({'success': False, 'error': 'Заголовок и текст обязательны'})
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('INSERT INTO site_news (title, body) VALUES (?, ?)', (title, body))
    news_id = cursor.lastrowid
    conn.commit()
    news = row_to_dict(conn.execute('SELECT * FROM site_news WHERE id = ?', (news_id,)).fetchone())
    conn.close()
    return jsonify({'success': True, 'news': news})


@app.route('/admin/news/delete/<int:news_id>', methods=['POST'])
def admin_delete_news(news_id):
    if 'user_id' not in session or session.get('role') != 'admin':
        return jsonify({'success': False, 'error': 'Нет доступа'}), 403
    conn = get_db_connection()
    conn.execute('DELETE FROM site_news WHERE id = ?', (news_id,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})


# ==================== ОСНОВНЫЕ МАРШРУТЫ ====================

@app.route('/')
def index():
    if 'user_id' in session:
        username = session.get('username')
        return redirect_based_on_role(username)
    return redirect('/login')


@app.route('/home', methods=['GET', 'POST'])
def home():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']

    if request.method == 'POST':
        # Создание нового поста
        content = request.form.get('content', '').strip()
        files = request.files.getlist('media_files')

        # Проверяем, есть ли хоть что-то для публикации
        has_content = bool(content)
        has_files = any(f and f.filename for f in files)

        if not has_content and not has_files:
            flash('Пост не может быть пустым. Добавьте текст или файлы.', 'error')
            return redirect('/home')

        conn = get_db_connection()
        try:
            current_datetime = get_current_datetime()

            # Создаем пост
            cursor = conn.cursor()
            cursor.execute('''
                INSERT INTO posts (user_id, content, created_at)
                VALUES (?, ?, ?)
            ''', (user_id, content, current_datetime))

            post_id = cursor.lastrowid
            print(f"Создан пост ID: {post_id} с текстом: {content}")

            # Обработка загруженных файлов
            uploaded_files = 0
            if files:
                for file in files:
                    if file and file.filename and file.filename.strip() != '':
                        original_name = file.filename
                        filename = secure_filename(file.filename)

                        if '.' not in filename:
                            continue

                        file_ext = filename.rsplit('.', 1)[1].lower()

                        # Определяем тип файла
                        if file_ext in ALLOWED_IMAGE_EXTENSIONS:
                            file_type = 'image'
                        elif file_ext in ALLOWED_VIDEO_EXTENSIONS:
                            file_type = 'video'
                        elif file_ext in ALLOWED_DOCUMENT_EXTENSIONS:
                            file_type = 'document'
                        else:
                            continue

                        # Генерируем уникальное имя файла
                        import time
                        unique_filename = f"post_{post_id}_{int(time.time())}_{hashlib.md5(filename.encode()).hexdigest()[:8]}.{file_ext}"
                        file_path = os.path.join(app.config['POST_MEDIA_FOLDER'], unique_filename)

                        try:
                            # Сохраняем файл
                            file.save(file_path)
                            print(f"Файл сохранен: {file_path}")

                            # Сохраняем в БД
                            cursor.execute('''
                                INSERT INTO post_media (post_id, filename, file_type, original_filename)
                                VALUES (?, ?, ?, ?)
                            ''', (post_id, unique_filename, file_type, original_name))

                            uploaded_files += 1
                            print(f"Сохранен файл {uploaded_files}: {unique_filename} (тип: {file_type})")

                        except Exception as file_error:
                            print(f"Ошибка при сохранении файла: {str(file_error)}")
                            continue

            conn.commit()

            if uploaded_files > 0:
                flash(f'Пост опубликован с {uploaded_files} файл(ов)!', 'success')
            else:
                flash('Пост опубликован!', 'success')

        except Exception as e:
            conn.rollback()
            print(f"Ошибка при создании поста: {e}")
            import traceback
            traceback.print_exc()
            flash(f'Ошибка при публикации поста: {str(e)[:50]}', 'error')
        finally:
            conn.close()

        return redirect('/home')

    # GET запрос - показываем ленту
    feed_items = []

    try:
        filter_type = request.args.get('filter', 'all')
        feed_items = get_posts_feed(user_id, limit=20, filter_type=filter_type, offset=0)

    except Exception as e:
        print(f"Ошибка при подготовке ленты: {e}")
        import traceback
        traceback.print_exc()
        flash('Не удалось загрузить все новости', 'info')

    # Получаем количество заявок в друзья
    friend_requests_count = 0
    conn = get_db_connection()
    try:
        result = conn.execute('''
            SELECT COUNT(*) as count FROM friendships 
            WHERE receiver_id = ? AND status = 'pending'
        ''', (user_id,)).fetchone()
        friend_requests_count = result['count'] if result else 0
    except Exception as e:
        print(f"Ошибка при получении заявок: {e}")
    finally:
        conn.close()

    conn3 = get_db_connection()
    site_news = rows_to_dicts(conn3.execute(
        'SELECT * FROM site_news ORDER BY created_at DESC LIMIT 10'
    ).fetchall())
    is_admin = session.get('role') == 'admin'
    conn3.close()

    return render_template('home.html',
                           username=session.get('username'),
                           friend_requests_count=friend_requests_count,
                           feed_items=feed_items,
                           site_news=site_news,
                           is_admin=is_admin)


@app.route('/profile')
def profile():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем основную информацию о пользователе
    user = row_to_dict(conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone())

    # Получаем профиль пользователя
    profile_data = row_to_dict(conn.execute('SELECT * FROM user_profiles WHERE user_id = ?', (user_id,)).fetchone())

    # Получаем количество постов
    posts_count = conn.execute('SELECT COUNT(*) as count FROM posts WHERE user_id = ?', (user_id,)).fetchone()['count']

    # Получаем количество друзей
    friends_count = conn.execute('''
        SELECT COUNT(*) as count FROM friendships 
        WHERE (sender_id = ? OR receiver_id = ?) AND status = 'accepted'
    ''', (user_id, user_id)).fetchone()['count']

    # Получаем количество заявок в друзья
    friend_requests_count = conn.execute('''
        SELECT COUNT(*) as count FROM friendships 
        WHERE receiver_id = ? AND status = 'pending'
    ''', (user_id,)).fetchone()['count']

    # Получаем количество пользователей в черном списке
    blacklist_count = conn.execute('''
        SELECT COUNT(*) as count FROM blacklist 
        WHERE blocker_id = ?
    ''', (user_id,)).fetchone()['count']

    conn.close()

    return render_template('profile.html',
                           user=user,
                           profile=profile_data,
                           posts_count=posts_count,
                           friends_count=friends_count,
                           friend_requests_count=friend_requests_count,
                           blacklist_count=blacklist_count)


@app.route('/profile/<int:user_id>')
def view_profile(user_id):
    if 'user_id' not in session:
        return redirect('/login')

    conn = get_db_connection()
    user = row_to_dict(conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone())

    if not user:
        flash('Пользователь не найден', 'error')
        return redirect('/home')

    profile_data = row_to_dict(conn.execute('SELECT * FROM user_profiles WHERE user_id = ?', (user_id,)).fetchone())
    posts_count = conn.execute('SELECT COUNT(*) as count FROM posts WHERE user_id = ?', (user_id,)).fetchone()['count']
    friend_status = None
    if user_id != session['user_id']:
        friendship = conn.execute('''
            SELECT status FROM friendships 
            WHERE (sender_id = ? AND receiver_id = ?) 
            OR (sender_id = ? AND receiver_id = ?)
        ''', (session['user_id'], user_id, user_id, session['user_id'])).fetchone()
        friend_status = friendship['status'] if friendship else None

    conn.close()

    return render_template('view_profile.html',
                           user=user,
                           profile=profile_data,
                           posts_count=posts_count,
                           friend_status=friend_status,
                           is_own_profile=(user_id == session['user_id']))


@app.route('/profile/edit', methods=['GET', 'POST'])
def edit_profile():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']

    if request.method == 'POST':
        full_name = request.form.get('full_name', '')
        bio = request.form.get('bio', '')
        location = request.form.get('location', '')

        avatar_filename = None
        if 'avatar' in request.files:
            file = request.files['avatar']
            if file and file.filename and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_ext = filename.rsplit('.', 1)[1].lower()
                unique_filename = f"{user_id}_{hashlib.md5(str(datetime.now()).encode()).hexdigest()[:10]}.{file_ext}"
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                file.save(file_path)
                avatar_filename = unique_filename

        conn = get_db_connection()
        existing_profile = conn.execute('SELECT id FROM user_profiles WHERE user_id = ?', (user_id,)).fetchone()

        if existing_profile:
            if avatar_filename:
                conn.execute('''
                    UPDATE user_profiles 
                    SET full_name = ?, bio = ?, location = ?, avatar = ?
                    WHERE user_id = ?
                ''', (full_name, bio, location, avatar_filename, user_id))
            else:
                conn.execute('''
                    UPDATE user_profiles 
                    SET full_name = ?, bio = ?, location = ?
                    WHERE user_id = ?
                ''', (full_name, bio, location, user_id))
        else:
            avatar_to_use = avatar_filename if avatar_filename else 'default_avatar.png'
            conn.execute('''
                INSERT INTO user_profiles (user_id, full_name, bio, location, avatar)
                VALUES (?, ?, ?, ?, ?)
            ''', (user_id, full_name, bio, location, avatar_to_use))

        conn.commit()
        conn.close()
        flash('Профиль успешно обновлен!', 'success')
        return redirect('/profile')

    # GET запрос - показываем форму редактирования
    conn = get_db_connection()
    profile_data = row_to_dict(conn.execute('SELECT * FROM user_profiles WHERE user_id = ?', (user_id,)).fetchone())
    conn.close()

    return render_template('edit_profile.html', profile=profile_data)


@app.route('/find_friends', methods=['GET', 'POST'])
def find_friends():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    search_results = []

    if request.method == 'POST':
        search_query = request.form.get('search_query', '').strip()
        if search_query:
            conn = get_db_connection()

            # Ищем пользователей по имени пользователя или ФИО
            rows = conn.execute('''
                SELECT u.id, u.username, up.full_name, 
                       COALESCE(up.avatar, 'default_avatar.png') as avatar 
                FROM users u
                LEFT JOIN user_profiles up ON u.id = up.user_id
                WHERE (u.username LIKE ? OR up.full_name LIKE ?) 
                AND u.id != ? AND u.is_banned = 0
                LIMIT 20
            ''', (f'%{search_query}%', f'%{search_query}%', user_id)).fetchall()

            # Преобразуем Row объекты в словари
            search_results = rows_to_dicts(rows)

            # Проверяем статус дружбы и черного списка для каждого найденного пользователя
            for user in search_results:
                # Проверяем статус дружбы
                friend_status = conn.execute('''
                    SELECT status FROM friendships 
                    WHERE (sender_id = ? AND receiver_id = ?) 
                    OR (sender_id = ? AND receiver_id = ?)
                ''', (user_id, user['id'], user['id'], user_id)).fetchone()

                user['friend_status'] = friend_status['status'] if friend_status else None

                # Проверяем, находится ли пользователь в черном списке
                blacklisted = conn.execute('''
                    SELECT id FROM blacklist 
                    WHERE blocker_id = ? AND blocked_id = ?
                ''', (user_id, user['id'])).fetchone()

                user['is_blacklisted'] = blacklisted is not None

            conn.close()

    return render_template('find_friends.html', search_results=search_results)


@app.route('/add_friend/<int:friend_id>', methods=['POST'])
def add_friend(friend_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']

    if user_id == friend_id:
        flash('Нельзя добавить себя в друзья!', 'error')
        return redirect('/find_friends')

    conn = get_db_connection()

    # Проверяем, не заблокирован ли пользователь
    is_blacklisted = conn.execute('''
        SELECT id FROM blacklist 
        WHERE blocker_id = ? AND blocked_id = ?
    ''', (user_id, friend_id)).fetchone()

    if is_blacklisted:
        flash('Вы не можете добавить в друзья пользователя из черного списка!', 'error')
        conn.close()
        return redirect('/find_friends')

    # Проверяем, не существует ли уже заявка
    existing_request = conn.execute('''
        SELECT * FROM friendships 
        WHERE (sender_id = ? AND receiver_id = ?) 
        OR (sender_id = ? AND receiver_id = ?)
    ''', (user_id, friend_id, friend_id, user_id)).fetchone()

    if existing_request:
        status = existing_request['status']
        if status == 'pending':
            if existing_request['receiver_id'] == user_id:
                # Пользователь принимает заявку
                conn.execute('''
                    UPDATE friendships SET status = 'accepted' 
                    WHERE id = ?
                ''', (existing_request['id'],))
                flash('Заявка в друзья принята!', 'success')
            else:
                flash('Заявка уже отправлена и ожидает подтверждения', 'info')
        elif status == 'accepted':
            flash('Этот пользователь уже у вас в друзьях', 'info')
        elif status == 'rejected':
            # Разрешаем отправить заявку заново
            conn.execute('''
                UPDATE friendships SET status = 'pending', sender_id = ?, receiver_id = ?
                WHERE id = ?
            ''', (user_id, friend_id, existing_request['id']))
            flash('Заявка в друзья отправлена повторно!', 'success')
    else:
        # Отправляем новую заявку
        conn.execute('''
            INSERT INTO friendships (sender_id, receiver_id, status)
            VALUES (?, ?, 'pending')
        ''', (user_id, friend_id))
        flash('Заявка в друзья отправлена!', 'success')

    conn.commit()
    conn.close()

    return redirect('/find_friends')


@app.route('/add_friend_to_blacklist/<int:friend_id>')
def add_friend_to_blacklist(friend_id):
    """Добавить друга в черный список из раздела друзей"""
    if 'user_id' not in session:
        return redirect('/login')

    blocker_id = session['user_id']

    if blocker_id == friend_id:
        flash('Нельзя добавить себя в черный список!', 'error')
        return redirect('/friends')

    conn = get_db_connection()

    # Получаем информацию о друге
    friend_info = conn.execute('''
        SELECT u.username, up.full_name 
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id = ?
    ''', (friend_id,)).fetchone()

    if not friend_info:
        conn.close()
        flash('Пользователь не найден', 'error')
        return redirect('/friends')

    # Проверяем, не добавлен ли уже пользователь в черный список
    existing = conn.execute('''
        SELECT id FROM blacklist 
        WHERE blocker_id = ? AND blocked_id = ?
    ''', (blocker_id, friend_id)).fetchone()

    if existing:
        flash('Пользователь уже в вашем черном списке', 'info')
    else:
        # Добавляем в черный список
        conn.execute('''
            INSERT INTO blacklist (blocker_id, blocked_id, reason)
            VALUES (?, ?, ?)
        ''', (blocker_id, friend_id, 'Добавлен из раздела друзей'))

        # Удаляем из друзей
        conn.execute('''
            DELETE FROM friendships 
            WHERE ((sender_id = ? AND receiver_id = ?) 
            OR (sender_id = ? AND receiver_id = ?)) 
            AND status = 'accepted'
        ''', (blocker_id, friend_id, friend_id, blocker_id))

        # Отменяем все заявки в друзья между этими пользователями
        conn.execute('''
            DELETE FROM friendships 
            WHERE (sender_id = ? AND receiver_id = ?) 
            OR (sender_id = ? AND receiver_id = ?)
        ''', (blocker_id, friend_id, friend_id, blocker_id))

        flash(f'{friend_info["username"]} добавлен в черный список и удален из друзей', 'success')

    conn.commit()
    conn.close()

    return redirect('/friends')


@app.route('/friends')
def friends():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем список друзей
    friends_rows = conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar 
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id IN (
            SELECT CASE 
                WHEN sender_id = ? THEN receiver_id 
                ELSE sender_id 
            END as friend_id
            FROM friendships 
            WHERE (sender_id = ? OR receiver_id = ?) AND status = 'accepted'
        )
        ORDER BY up.full_name, u.username
    ''', (user_id, user_id, user_id)).fetchall()

    # Преобразуем в словари
    friends = rows_to_dicts(friends_rows)

    # Получаем входящие заявки в друзья
    incoming_requests_rows = conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar, 
               f.created_at, f.id as request_id
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        JOIN friendships f ON u.id = f.sender_id
        WHERE f.receiver_id = ? AND f.status = 'pending'
        ORDER BY f.created_at DESC
    ''', (user_id,)).fetchall()

    incoming_requests = rows_to_dicts(incoming_requests_rows)

    # Получаем исходящие заявки
    outgoing_requests_rows = conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar, 
               f.created_at, f.id as request_id
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        JOIN friendships f ON u.id = f.receiver_id
        WHERE f.sender_id = ? AND f.status = 'pending'
        ORDER BY f.created_at DESC
    ''', (user_id,)).fetchall()

    outgoing_requests = rows_to_dicts(outgoing_requests_rows)

    conn.close()

    conn.close()

    return render_template('friends.html',
                           friends=friends,
                           incoming_requests=incoming_requests,
                           outgoing_requests=outgoing_requests)


@app.route('/friends/pending_count')
def friends_pending_count():
    """Количество входящих заявок в друзья (для бейджа)"""
    if 'user_id' not in session:
        return jsonify({'count': 0})
    conn = get_db_connection()
    count = conn.execute('''
        SELECT COUNT(*) as count FROM friendships
        WHERE receiver_id = ? AND status = 'pending'
    ''', (session['user_id'],)).fetchone()['count']
    conn.close()
    return jsonify({'count': count})


@app.route('/friend_action/<int:request_id>/<action>')
def friend_action(request_id, action):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем информацию о заявке
    friend_request = conn.execute('''
        SELECT * FROM friendships 
        WHERE id = ?
    ''', (request_id,)).fetchone()

    if friend_request:
        friend_request_dict = dict(friend_request)
        # Проверяем, имеет ли пользователь право выполнять это действие
        if action == 'accept' and friend_request_dict['receiver_id'] == user_id:
            conn.execute('''
                UPDATE friendships SET status = 'accepted' 
                WHERE id = ?
            ''', (request_id,))
            flash('Заявка в друзья принята!', 'success')
        elif action == 'reject' and friend_request_dict['receiver_id'] == user_id:
            conn.execute('''
                UPDATE friendships SET status = 'rejected' 
                WHERE id = ?
            ''', (request_id,))
            flash('Заявка в друзья отклонена', 'info')
        elif action == 'cancel' and friend_request_dict['sender_id'] == user_id:
            # Пользователь отменяет свою исходящую заявку
            conn.execute('''
                DELETE FROM friendships WHERE id = ?
            ''', (request_id,))
            flash('Заявка отменена', 'info')
        else:
            flash('У вас нет прав для выполнения этого действия', 'error')

    conn.commit()
    conn.close()

    return redirect('/friends')


@app.route('/remove_friend/<int:friend_id>')
def remove_friend(friend_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Удаляем запись о дружбе (любой статус)
    conn.execute('''
        DELETE FROM friendships 
        WHERE (sender_id = ? AND receiver_id = ?) 
        OR (sender_id = ? AND receiver_id = ?)
    ''', (user_id, friend_id, friend_id, user_id))

    conn.commit()
    conn.close()

    flash('Пользователь удален из друзей', 'info')
    return redirect('/friends')


# Черный список
@app.route('/blacklist')
def blacklist():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем список пользователей в черном списке
    blacklist_rows = conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar,
               b.reason, b.created_at
        FROM blacklist b
        JOIN users u ON b.blocked_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE b.blocker_id = ?
        ORDER BY b.created_at DESC
    ''', (user_id,)).fetchall()

    blacklisted_users = rows_to_dicts(blacklist_rows)

    conn.close()

    return render_template('blacklist.html', blacklisted_users=blacklisted_users)


@app.route('/add_to_blacklist/<int:user_id>', methods=['GET', 'POST'])
def add_to_blacklist(user_id):
    if 'user_id' not in session:
        return redirect('/login')

    blocker_id = session['user_id']

    if blocker_id == user_id:
        flash('Нельзя добавить себя в черный список!', 'error')
        return redirect(request.referrer or '/find_friends')

    if request.method == 'POST':
        reason = request.form.get('reason', '')

        conn = get_db_connection()

        # Проверяем, не добавлен ли уже пользователь в черный список
        existing = conn.execute('''
            SELECT id FROM blacklist 
            WHERE blocker_id = ? AND blocked_id = ?
        ''', (blocker_id, user_id)).fetchone()

        if existing:
            flash('Пользователь уже в вашем черном списке', 'info')
        else:
            # Добавляем в черный список
            conn.execute('''
                INSERT INTO blacklist (blocker_id, blocked_id, reason)
                VALUES (?, ?, ?)
            ''', (blocker_id, user_id, reason))

            # Удаляем из друзей, если были друзьями
            conn.execute('''
                DELETE FROM friendships 
                WHERE ((sender_id = ? AND receiver_id = ?) 
                OR (sender_id = ? AND receiver_id = ?)) 
                AND status = 'accepted'
            ''', (blocker_id, user_id, user_id, blocker_id))

            # Отменяем все заявки в друзья между этими пользователями
            conn.execute('''
                DELETE FROM friendships 
                WHERE (sender_id = ? AND receiver_id = ?) 
                OR (sender_id = ? AND receiver_id = ?)
            ''', (blocker_id, user_id, user_id, blocker_id))

            flash('Пользователь добавлен в черный список', 'success')

        conn.commit()
        conn.close()

        return redirect('/blacklist')

    # GET запрос - показываем форму
    conn = get_db_connection()
    user_info = conn.execute('''
        SELECT u.username, up.full_name 
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id = ?
    ''', (user_id,)).fetchone()

    conn.close()

    if user_info:
        return render_template('add_to_blacklist.html',
                               user_id=user_id,
                               username=user_info['username'],
                               full_name=user_info['full_name'])
    else:
        flash('Пользователь не найден', 'error')
        return redirect('/find_friends')


@app.route('/remove_from_blacklist/<int:blocked_id>')
def remove_from_blacklist(blocked_id):
    if 'user_id' not in session:
        return redirect('/login')

    blocker_id = session['user_id']

    conn = get_db_connection()

    # Удаляем из черного списка
    conn.execute('''
        DELETE FROM blacklist 
        WHERE blocker_id = ? AND blocked_id = ?
    ''', (blocker_id, blocked_id))

    conn.commit()
    conn.close()

    flash('Пользователь удален из черного списка', 'info')
    return redirect('/blacklist')


# Жалобы на пользователей
@app.route('/report_friend/<int:friend_id>', methods=['GET', 'POST'])
def report_friend(friend_id):
    """Отправить жалобу на друга из раздела друзей"""
    if 'user_id' not in session:
        return redirect('/login')

    reporter_id = session['user_id']

    if reporter_id == friend_id:
        flash('Нельзя отправить жалобу на себя!', 'error')
        return redirect('/friends')

    if request.method == 'POST':
        reason = request.form.get('reason', '').strip()
        violation_type = request.form.get('violation_type', 'other')

        if not reason or len(reason) < 10:
            flash('Пожалуйста, опишите причину жалобы подробнее (минимум 10 символов)', 'error')
            return redirect(f'/report_friend/{friend_id}')

        conn = get_db_connection()

        try:
            # Отправляем жалобу
            full_reason = f"Тип нарушения: {violation_type}\n\n{reason}"
            conn.execute('''
                INSERT INTO reports (reporter_id, reported_id, reason, status)
                VALUES (?, ?, ?, 'pending')
            ''', (reporter_id, friend_id, full_reason))

            conn.commit()
            flash('Жалоба отправлена администратору. Спасибо за ваше сообщение!', 'success')

        except Exception as e:
            conn.rollback()
            print(f"Ошибка при сохранении жалобы: {e}")
            flash('Ошибка при отправке жалобы. Попробуйте позже.', 'error')
        finally:
            conn.close()

        return redirect('/friends')

    # GET запрос - показываем форму
    conn = get_db_connection()
    friend_info = conn.execute('''
        SELECT u.username, up.full_name 
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id = ?
    ''', (friend_id,)).fetchone()

    conn.close()

    if friend_info:
        return render_template('report_user.html',
                               user_id=friend_id,
                               username=friend_info['username'],
                               full_name=friend_info['full_name'],
                               source='friends')
    else:
        flash('Пользователь не найден', 'error')
        return redirect('/friends')


@app.route('/report_user/<int:user_id>', methods=['GET', 'POST'])
def report_user(user_id):
    """Отправить жалобу на пользователя из поиска"""
    if 'user_id' not in session:
        return redirect('/login')

    reporter_id = session['user_id']

    if reporter_id == user_id:
        flash('Нельзя отправить жалобу на себя!', 'error')
        return redirect('/find_friends')

    if request.method == 'POST':
        reason = request.form.get('reason', '').strip()
        violation_type = request.form.get('violation_type', 'other')

        if not reason or len(reason) < 10:
            flash('Пожалуйста, опишите причину жалобы подробнее (минимум 10 символов)', 'error')
            return redirect(f'/report_user/{user_id}')

        conn = get_db_connection()

        try:
            # Отправляем жалобу
            full_reason = f"Тип нарушения: {violation_type}\n\n{reason}"
            conn.execute('''
                INSERT INTO reports (reporter_id, reported_id, reason, status)
                VALUES (?, ?, ?, 'pending')
            ''', (reporter_id, user_id, full_reason))

            conn.commit()
            flash('Жалоба отправлена администратору. Спасибо за ваше сообщение!', 'success')

        except Exception as e:
            conn.rollback()
            print(f"Ошибка при сохранении жалобы: {e}")
            flash('Ошибка при отправке жалобы. Попробуйте позже.', 'error')
        finally:
            conn.close()

        return redirect('/find_friends')

    # GET запрос - показываем форму
    conn = get_db_connection()
    user_info = conn.execute('''
        SELECT u.username, up.full_name 
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id = ?
    ''', (user_id,)).fetchone()

    conn.close()
    if user_info:
        return render_template('report_user.html',
                               user_id=user_id,
                               username=user_info['username'],
                               full_name=user_info['full_name'],
                               source='find_friends')
    else:
        flash('Пользователь не найден', 'error')
        return redirect('/find_friends')


# ==================== ГРУППЫ ====================

@app.route('/groups')
def groups_list():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    search_query = request.args.get('search', '')
    tab = request.args.get('tab', 'all')

    conn = get_db_connection()

    # Базовый запрос для ВСЕХ групп
    query = '''
        SELECT g.*, 
               COUNT(DISTINCT gm.user_id) as members_count,
               COUNT(DISTINCT gp.id) as posts_count,
               EXISTS(SELECT 1 FROM group_members WHERE group_id = g.id AND user_id = ?) as is_member,
               (SELECT role FROM group_members WHERE group_id = g.id AND user_id = ? LIMIT 1) as role,
               EXISTS(SELECT 1 FROM group_requests WHERE group_id = g.id AND user_id = ? AND status = 'pending') as has_pending_request
        FROM groups g
        LEFT JOIN group_members gm ON g.id = gm.group_id
        LEFT JOIN group_posts gp ON g.id = gp.group_id
    '''
    params = [user_id, user_id, user_id]

    if tab == 'my':
        query += ' WHERE g.id IN (SELECT group_id FROM group_members WHERE user_id = ?)'
        params.append(user_id)
    # Для вкладки "Все группы" показываем ВСЕ группы (и публичные, и приватные)

    if search_query:
        if 'WHERE' in query:
            query += ' AND'
        else:
            query += ' WHERE'
        query += ' (g.name LIKE ? OR g.description LIKE ?)'
        params.extend([f'%{search_query}%', f'%{search_query}%'])

    query += ' GROUP BY g.id ORDER BY members_count DESC'

    groups = rows_to_dicts(conn.execute(query, params).fetchall())

    # Получаем creator_id для каждой группы
    for group in groups:
        creator = conn.execute('SELECT creator_id FROM groups WHERE id = ?', (group['id'],)).fetchone()
        group['creator_id'] = creator['creator_id'] if creator else None

    conn.close()

    return render_template('groups.html',
                           groups=groups,
                           search_query=search_query,
                           current_tab=tab,
                           session_user_id=user_id)


@app.route('/create_group', methods=['POST'])
def create_group():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    name = request.form.get('name', '').strip()
    description = request.form.get('description', '').strip()
    is_public = request.form.get('is_public') == 'on'

    if not name:
        flash('Название паблика обязательно', 'error')
        return redirect('/groups')

    avatar_filename = None
    if 'avatar' in request.files:
        file = request.files['avatar']
        if file and file.filename and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_ext = filename.rsplit('.', 1)[1].lower()
            unique_filename = f"group_{hashlib.md5(str(datetime.now()).encode()).hexdigest()[:10]}.{file_ext}"
            file_path = os.path.join(app.config['GROUP_UPLOAD_FOLDER'], unique_filename)
            file.save(file_path)
            avatar_filename = unique_filename

    conn = get_db_connection()
    current_datetime = get_current_datetime()

    try:
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO groups (name, description, avatar, creator_id, is_public, created_at)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (name, description, avatar_filename if avatar_filename else 'default_group.png',
              user_id, 1 if is_public else 0, current_datetime))

        group_id = cursor.lastrowid

        cursor.execute('''
            INSERT INTO group_members (group_id, user_id, role)
            VALUES (?, ?, 'admin')
        ''', (group_id, user_id))

        conn.commit()
        flash('Паблик успешно создан!', 'success')
        return redirect(f'/group/{group_id}')

    except Exception as e:
        conn.rollback()
        flash(f'Ошибка при создании паблика: {str(e)}', 'error')
        return redirect('/groups')
    finally:
        conn.close()


@app.route('/join_group/<int:group_id>', methods=['POST'])
def join_group(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    try:
        existing_member = conn.execute('''
            SELECT id FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, user_id)).fetchone()

        if existing_member:
            flash('Вы уже в этой группе', 'info')
            return redirect(request.referrer or f'/group/{group_id}')

        group = conn.execute('SELECT * FROM groups WHERE id = ?', (group_id,)).fetchone()
        if not group:
            flash('Группа не найдена', 'error')
            return redirect('/groups')

        if group['is_public']:
            conn.execute('''
                INSERT INTO group_members (group_id, user_id, role)
                VALUES (?, ?, 'member')
            ''', (group_id, user_id))
            flash(f'Вы присоединились к группе "{group["name"]}"!', 'success')
        else:
            existing_request = conn.execute('''
                SELECT id FROM group_requests 
                WHERE group_id = ? AND user_id = ? AND status = 'pending'
            ''', (group_id, user_id)).fetchone()

            if existing_request:
                flash('Вы уже отправили заявку на вступление', 'info')
            else:
                conn.execute('''
                    INSERT INTO group_requests (group_id, user_id, status, created_at)
                    VALUES (?, ?, 'pending', ?)
                ''', (group_id, user_id, get_current_datetime()))
                flash('Заявка на вступление отправлена! Администратор рассмотрит её в ближайшее время.', 'success')

        conn.commit()

    except Exception as e:
        conn.rollback()
        flash(f'Ошибка: {str(e)}', 'error')
    finally:
        conn.close()

    return redirect(request.referrer or f'/group/{group_id}')


@app.route('/leave_group/<int:group_id>', methods=['POST'])
def leave_group(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    try:
        group = conn.execute('SELECT creator_id, name FROM groups WHERE id = ?', (group_id,)).fetchone()
        membership = conn.execute('SELECT role FROM group_members WHERE group_id = ? AND user_id = ?',
                                  (group_id, user_id)).fetchone()

        # Если пользователь — администратор группы, удаляем группу целиком
        if membership and membership['role'] == 'admin':
            conn.execute('DELETE FROM group_posts WHERE group_id = ?', (group_id,))
            conn.execute(
                'DELETE FROM group_post_likes WHERE post_id IN (SELECT id FROM group_posts WHERE group_id = ?)',
                (group_id,))
            conn.execute('DELETE FROM group_members WHERE group_id = ?', (group_id,))
            conn.execute('DELETE FROM group_requests WHERE group_id = ?', (group_id,))
            conn.execute('DELETE FROM groups WHERE id = ?', (group_id,))

            group_avatar = conn.execute('SELECT avatar FROM groups WHERE id = ?', (group_id,)).fetchone()
            if group_avatar and group_avatar['avatar'] and group_avatar['avatar'] != 'default_group.png':
                try:
                    os.remove(os.path.join(app.config['GROUP_UPLOAD_FOLDER'], group_avatar['avatar']))
                except:
                    pass

            flash(f'Группа "{group["name"]}" была удалена, так как вы покинули её как создатель.', 'info')

        else:
            conn.execute('DELETE FROM group_members WHERE group_id = ? AND user_id = ?', (group_id, user_id))
            flash(f'Вы отписались от группы "{group["name"]}"', 'info')

        conn.commit()

    except Exception as e:
        conn.rollback()
        flash(f'Ошибка: {str(e)}', 'error')

    finally:
        conn.close()

    return redirect('/groups')


@app.route('/group/<int:group_id>')
def group_detail(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем информацию о группе
    group = row_to_dict(conn.execute('''
        SELECT g.*, 
               COUNT(DISTINCT gm.user_id) as members_count,
               COUNT(DISTINCT gp.id) as posts_count,
               u.username as creator_username
        FROM groups g
        LEFT JOIN group_members gm ON g.id = gm.group_id
        LEFT JOIN group_posts gp ON g.id = gp.group_id
        LEFT JOIN users u ON g.creator_id = u.id
        WHERE g.id = ?
        GROUP BY g.id
    ''', (group_id,)).fetchone())

    if not group:
        flash('Паблик не найден', 'error')
        return redirect('/groups')

    # Проверяем, является ли пользователь участником
    membership = conn.execute('''
        SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
    ''', (group_id, user_id)).fetchone()

    is_member = membership is not None
    role = membership['role'] if membership else None

    # Проверяем, может ли пользователь создавать посты
    can_post = False
    if is_member:
        can_post = check_post_permission(group_id, user_id, conn)

    # Проверяем, показывать ли вкладку "Заявки"
    show_requests_tab = False
    pending_requests = []
    pending_requests_count = 0

    if is_member and role in ['admin', 'moderator'] and not group['is_public']:
        # Проверяем права на просмотр заявок
        if group.get('request_permissions') == 'admins':
            show_requests_tab = role == 'admin'
        else:  # moderators или по умолчанию
            show_requests_tab = role in ['admin', 'moderator']

        if show_requests_tab:
            # Получаем заявки на вступление
            pending_requests = rows_to_dicts(conn.execute('''
                SELECT gr.id as request_id, gr.user_id, gr.created_at,
                       u.username, up.full_name, COALESCE(up.avatar, 'default_avatar.png') as avatar
                FROM group_requests gr
                JOIN users u ON gr.user_id = u.id
                LEFT JOIN user_profiles up ON u.id = up.user_id
                WHERE gr.group_id = ? AND gr.status = 'pending'
                ORDER BY gr.created_at
            ''', (group_id,)).fetchall())

            pending_requests_count = len(pending_requests)

    # Получаем посты группы
    posts_data = []
    if is_member or group['is_public']:
        posts = rows_to_dicts(conn.execute('''
            SELECT gp.*,
                   u.username as author_username,
                   up.full_name as author_name,
                   COALESCE(up.avatar, 'default_avatar.png') as author_avatar,
                   (SELECT COUNT(*) FROM group_post_likes WHERE post_id = gp.id) as likes_count
            FROM group_posts gp
            JOIN users u ON gp.author_id = u.id
            LEFT JOIN user_profiles up ON u.id = up.user_id
            WHERE gp.group_id = ?
            ORDER BY gp.created_at DESC
            LIMIT 50
        ''', (group_id,)).fetchall())

        # Добавляем роль автора поста для каждого поста
        for post in posts:
            # Получаем роль автора поста в этой группе
            author_role = conn.execute('''
                SELECT role FROM group_members 
                WHERE group_id = ? AND user_id = ?
            ''', (group_id, post['author_id'])).fetchone()

            # Если автор не состоит в группе (хотя это маловероятно), считаем его участником
            post['author_role'] = author_role['role'] if author_role else 'member'

            # Получаем медиафайлы для поста
            post['media_files'] = get_post_media(post['id'], is_group_post=True)

            posts_data.append(post)

    # Получаем участников группы
    members = rows_to_dicts(conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar,


        gm.role, gm.joined_at
        FROM group_members gm
        JOIN users u ON gm.user_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE gm.group_id = ?
        ORDER BY 
            CASE gm.role 
                WHEN 'admin' THEN 1
                WHEN 'moderator' THEN 2
                ELSE 3
            END,
            u.username
    ''', (group_id,)).fetchall())

    # Получаем создателя
    creator = row_to_dict(conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id = ?
    ''', (group['creator_id'],)).fetchone())

    # Проверяем, есть ли у пользователя ожидающая заявка (для приватных групп)
    has_pending_request = False
    if not is_member and not group['is_public']:
        pending_request = conn.execute('''
            SELECT id FROM group_requests 
            WHERE group_id = ? AND user_id = ? AND status = 'pending'
        ''', (group_id, user_id)).fetchone()
        has_pending_request = pending_request is not None

    conn.close()

    return render_template('group_detail.html',
                           group=group,
                           posts=posts_data,
                           members=members,
                           creator=creator,
                           is_member=is_member,
                           role=role,
                           can_post=can_post,
                           has_pending_request=has_pending_request,
                           session_user_id=user_id,
                           show_requests_tab=show_requests_tab,
                           pending_requests=pending_requests,
                           pending_requests_count=pending_requests_count)


@app.route('/group/<int:group_id>/create_post', methods=['POST'])
def create_group_post(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    content = request.form.get('content', '').strip()
    conn = get_db_connection()

    print(f"=== НАЧАЛО СОЗДАНИЯ ПОСТА ===")
    print(f"user_id: {user_id}")
    print(f"content: '{content}'")

    try:
        # Проверяем права на публикацию
        if not check_post_permission(group_id, user_id, conn):
            flash('У вас нет прав для создания постов в этой группе', 'error')
            return redirect(f'/group/{group_id}')

        # Проверяем, есть ли файлы
        has_valid_files = False
        files = []
        if 'media_files' in request.files:
            files = request.files.getlist('media_files')
            print(f"Получено файлов: {len(files)}")
            for file in files:
                if file and file.filename and file.filename.strip() != '':
                    has_valid_files = True

        # Проверяем, есть ли хоть что-то для публикации
        if not content and not has_valid_files:
            flash('Пост не может быть пустым. Добавьте текст или файлы.', 'error')
            return redirect(f'/group/{group_id}')

        current_datetime = get_current_datetime()

        # Создаем пост
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO group_posts (group_id, author_id, content, created_at)
            VALUES (?, ?, ?, ?)
        ''', (group_id, user_id, content, current_datetime))

        post_id = cursor.lastrowid
        print(f"Создан пост ID: {post_id}")

        # Обработка загруженных файлов
        uploaded_files = 0
        if files:
            for file in files:
                if file and file.filename and file.filename.strip() != '':
                    original_name = file.filename
                    filename = secure_filename(file.filename)

                    if '.' not in filename:
                        continue

                    file_ext = filename.rsplit('.', 1)[1].lower()

                    # Определяем тип файла
                    if file_ext in ALLOWED_IMAGE_EXTENSIONS:
                        file_type = 'image'
                    elif file_ext in ALLOWED_VIDEO_EXTENSIONS:
                        file_type = 'video'
                    elif file_ext in ALLOWED_DOCUMENT_EXTENSIONS:
                        file_type = 'document'
                    else:
                        continue

                    # Генерируем уникальное имя файла
                    import time
                    unique_filename = f"{post_id}_{int(time.time())}_{hashlib.md5(filename.encode()).hexdigest()[:8]}.{file_ext}"
                    file_path = os.path.join(app.config['POST_MEDIA_FOLDER'], unique_filename)

                    try:
                        # Сохраняем файл
                        file.save(file_path)

                        # ВАЖНО: Используем group_post_id для групповых постов
                        cursor.execute('''
                            INSERT INTO post_media (group_post_id, filename, file_type, original_filename)
                            VALUES (?, ?, ?, ?)
                        ''', (post_id, unique_filename, file_type, original_name))

                        uploaded_files += 1

                    except Exception as file_error:
                        print(f"Ошибка при сохранении файла: {str(file_error)}")
                        continue

        conn.commit()

        if uploaded_files > 0:
            flash(f'Пост опубликован в группе! Загружено {uploaded_files} файл(ов)', 'success')
        else:
            flash('Пост опубликован в группе!', 'success')

    except Exception as e:
        conn.rollback()
        print(f"Ошибка: {str(e)}")
        flash(f'Ошибка при создании поста: {str(e)}', 'error')
    finally:
        conn.close()

    return redirect(f'/group/{group_id}')


@app.route('/group/<int:group_id>/manage')
def manage_group(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']

    conn = get_db_connection()

    membership = conn.execute(''' SELECT role FROM group_members WHERE group_id = ? AND user_id = ? ''',
                              (group_id, user_id)).fetchone()

    if not membership or membership['role'] != 'admin':
        flash('У вас нет прав для управления этой группой', 'error')
        return redirect(f'/group/{group_id}')

    group = row_to_dict(conn.execute('SELECT * FROM groups WHERE id = ?', (group_id,)).fetchone())

    members = rows_to_dicts(conn.execute('''
        SELECT u.id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar,
               gm.role, gm.joined_at
        FROM group_members gm
        JOIN users u ON gm.user_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE gm.group_id = ?
        ORDER BY 
            CASE gm.role 
                WHEN 'admin' THEN 1
                WHEN 'moderator' THEN 2
                ELSE 3
            END,
            u.username
    ''', (group_id,)).fetchall())

    conn.close()

    return render_template('manage_group.html', group=group, members=members)


@app.route('/group/<int:group_id>/settings')
def group_settings(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Проверяем права (только админы)
    membership = conn.execute('''
        SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
    ''', (group_id, user_id)).fetchone()

    if not membership or membership['role'] != 'admin':
        flash('Только администраторы могут изменять настройки группы', 'error')
        return redirect(f'/group/{group_id}')

    # Получаем информацию о группе
    group = row_to_dict(conn.execute('SELECT * FROM groups WHERE id = ?', (group_id,)).fetchone())

    # Получаем список всех участников группы
    members = rows_to_dicts(conn.execute('''
        SELECT u.id as user_id, u.username, up.full_name, 
               COALESCE(up.avatar, 'default_avatar.png') as avatar,
               gm.role, gm.joined_at
        FROM group_members gm
        JOIN users u ON gm.user_id = u.id
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE gm.group_id = ?
        ORDER BY 
            CASE gm.role 
                WHEN 'admin' THEN 1
                WHEN 'moderator' THEN 2
                ELSE 3
            END,
            u.username
    ''', (group_id,)).fetchall())

    conn.close()
    return render_template('group_settings.html',
                           group=group,
                           members=members,
                           session_user_id=user_id)


@app.route('/group/<int:group_id>/settings', methods=['POST'])
def update_group_settings(group_id):
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Проверяем права
    membership = conn.execute('''
        SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
    ''', (group_id, user_id)).fetchone()

    if not membership or membership['role'] != 'admin':
        flash('Только администраторы могут изменять настройки группы', 'error')
        return redirect(f'/group/{group_id}')

    # Получаем данные из формы
    name = request.form.get('name', '').strip()
    description = request.form.get('description', '').strip()
    is_public = request.form.get('is_public') == '1'
    post_permissions = request.form.get('post_permissions', 'all')
    request_permissions = request.form.get('request_permissions', 'moderators')

    if not name:
        flash('Название группы обязательно', 'error')
        return redirect(f'/group/{group_id}/settings')

    try:
        # Обновляем основные данные
        conn.execute('''
            UPDATE groups 
            SET name = ?, description = ?, is_public = ?, 
                post_permissions = ?, request_permissions = ?
            WHERE id = ?
        ''', (name, description, 1 if is_public else 0,
              post_permissions, request_permissions, group_id))

        # Обработка аватарки
        if 'avatar' in request.files:
            file = request.files['avatar']
            if file and file.filename and allowed_file(file.filename):
                # Удаляем старый аватар если нужно
                old_avatar = conn.execute('SELECT avatar FROM groups WHERE id = ?', (group_id,)).fetchone()['avatar']

                filename = secure_filename(file.filename)
                file_ext = filename.rsplit('.', 1)[1].lower()
                unique_filename = f"group_{group_id}_{hashlib.md5(str(datetime.now()).encode()).hexdigest()[:8]}.{file_ext}"
                file_path = os.path.join(app.config['GROUP_UPLOAD_FOLDER'], unique_filename)
                file.save(file_path)

                # Обновляем в БД
                conn.execute('UPDATE groups SET avatar = ? WHERE id = ?', (unique_filename, group_id))

                # Удаляем старый файл (кроме дефолтного)
                if old_avatar and old_avatar != 'default_group.png':
                    try:
                        os.remove(os.path.join(app.config['GROUP_UPLOAD_FOLDER'], old_avatar))
                    except:
                        pass

        # Удаление аватарки
        if request.form.get('remove_avatar'):
            old_avatar = conn.execute('SELECT avatar FROM groups WHERE id = ?', (group_id,)).fetchone()['avatar']
            if old_avatar and old_avatar != 'default_group.png':
                try:
                    os.remove(os.path.join(app.config['GROUP_UPLOAD_FOLDER'], old_avatar))
                except:
                    pass
            conn.execute('UPDATE groups SET avatar = "default_group.png" WHERE id = ?', (group_id,))

        conn.commit()
        flash('Настройки группы обновлены!', 'success')

    except Exception as e:
        conn.rollback()
        flash(f'Ошибка при обновлении настроек: {str(e)}', 'error')
        return redirect(f'/group/{group_id}/settings')
    finally:
        conn.close()

    return redirect(f'/group/{group_id}')


@app.route('/group/<int:group_id>/request/<int:request_id>/<action>')
def handle_group_request(group_id, request_id, action):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем настройки группы
    group = conn.execute('SELECT request_permissions FROM groups WHERE id = ?', (group_id,)).fetchone()

    if not group:
        return jsonify({'success': False, 'error': 'Группа не найдена'})

    # Проверяем права в зависимости от настроек
    membership = conn.execute('''
        SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
    ''', (group_id, user_id)).fetchone()

    if not membership:
        return jsonify({'success': False, 'error': 'Вы не состоите в этой группе'})

    if group['request_permissions'] == 'admins':
        # Только админы могут обрабатывать заявки
        if membership['role'] != 'admin':
            return jsonify({'success': False, 'error': 'Только администраторы могут обрабатывать заявки'})
    else:
        # Админы и модераторы могут обрабатывать заявки
        if membership['role'] not in ['admin', 'moderator']:
            return jsonify({'success': False, 'error': 'Только администраторы и модераторы могут обрабатывать заявки'})

    try:
        if action == 'approve':
            # Получаем заявку
            request_data = conn.execute('SELECT * FROM group_requests WHERE id = ?', (request_id,)).fetchone()
            if request_data:
                # Добавляем пользователя в группу
                conn.execute('''
                    INSERT INTO group_members (group_id, user_id, role)
                    VALUES (?, ?, 'member')
                ''', (group_id, request_data['user_id']))

                # Удаляем заявку
                conn.execute('DELETE FROM group_requests WHERE id = ?', (request_id,))

                conn.commit()
                return jsonify({'success': True})


        elif action == 'reject':
            conn.execute('DELETE FROM group_requests WHERE id = ?', (request_id,))
            conn.commit()
            return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/group/<int:group_id>/change_role/<int:user_id>', methods=['POST'])
def change_member_role(group_id, user_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    current_user_id = session['user_id']
    data = request.get_json()
    new_role = data.get('role')

    if not new_role or new_role not in ['admin', 'moderator', 'member']:
        return jsonify({'success': False, 'error': 'Invalid role'})

    conn = get_db_connection()

    try:
        current_user_role = conn.execute('''
            SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, current_user_id)).fetchone()

        if not current_user_role or current_user_role['role'] != 'admin':
            return jsonify({'success': False, 'error': 'Только администраторы могут изменять роли'})

        target_role = conn.execute('''
            SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, user_id)).fetchone()

        if not target_role:
            return jsonify({'success': False, 'error': 'Пользователь не найден в группе'})

        if target_role['role'] == 'admin':
            return jsonify({'success': False, 'error': 'Нельзя изменить роль администратора — сначала передайте права'})

        conn.execute('''
            UPDATE group_members 
            SET role = ?
            WHERE group_id = ? AND user_id = ?
        ''', (new_role, group_id, user_id))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/group/<int:group_id>/transfer_admin/<int:new_admin_id>', methods=['POST'])
def transfer_admin_rights(group_id, new_admin_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    current_user_id = session['user_id']

    conn = get_db_connection()

    try:
        current_user_role = conn.execute('''
            SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, current_user_id)).fetchone()

        if not current_user_role or current_user_role['role'] != 'admin':
            return jsonify({'success': False, 'error': 'Только администраторы могут передавать права'})

        new_admin = conn.execute('''
            SELECT user_id, role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, new_admin_id)).fetchone()

        if not new_admin:
            return jsonify({'success': False, 'error': 'Пользователь не найден в группе'})

        if new_admin['role'] == 'admin':
            return jsonify({'success': False, 'error': 'Этот пользователь уже является администратором'})

        conn.execute('''
            UPDATE group_members 
            SET role = 'member'
            WHERE group_id = ? AND user_id = ?
        ''', (group_id, current_user_id))

        conn.execute('''
            UPDATE group_members 
            SET role = 'admin'
            WHERE group_id = ? AND user_id = ?
        ''', (group_id, new_admin_id))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/group/<int:group_id>/remove_member/<int:user_id>', methods=['POST'])
def remove_group_member(group_id, user_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    current_user_id = session['user_id']

    conn = get_db_connection()

    try:
        current_user_role = conn.execute('''
            SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, current_user_id)).fetchone()

        if not current_user_role or current_user_role['role'] != 'admin':
            return jsonify({'success': False, 'error': 'Только администраторы могут удалять участников'})

        target_user_role = conn.execute('''
            SELECT role FROM group_members WHERE group_id = ? AND user_id = ?
        ''', (group_id, user_id)).fetchone()

        if not target_user_role:
            return jsonify({'success': False, 'error': 'Пользователь не найден в группе'})

        if target_user_role['role'] == 'admin':
            return jsonify({'success': False, 'error': 'Нельзя исключить администратора — сначала передайте права'})

        conn.execute('DELETE FROM group_members WHERE group_id = ? AND user_id = ?', (group_id, user_id))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/group_post/edit/<int:post_id>', methods=['POST'])
def edit_group_post(post_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    user_id = session['user_id']
    data = request.get_json()
    new_content = data.get('content', '').strip()

    if not new_content:
        return jsonify({'success': False, 'error': 'Пост не может быть пустым'})

    conn = get_db_connection()

    try:
        post = conn.execute('''
            SELECT gp.*, g.creator_id, gp.author_id as post_author_id
            FROM group_posts gp
            JOIN groups g ON gp.group_id = g.id
            WHERE gp.id = ?
        ''', (post_id,)).fetchone()

        if not post:
            return jsonify({'success': False, 'error': 'Пост не найден'})

        user_role = conn.execute('''
            SELECT role FROM group_members 
            WHERE group_id = ? AND user_id = ?
        ''', (post['group_id'], user_id)).fetchone()

        author_role = conn.execute('''
            SELECT role FROM group_members 
            WHERE group_id = ? AND user_id = ?
        ''', (post['group_id'], post['post_author_id'])).fetchone()

        can_edit = False

        if post['post_author_id'] == user_id:
            can_edit = True
        elif user_role and user_role['role'] == 'admin':
            can_edit = True
        elif user_role and user_role['role'] == 'moderator':
            if author_role and author_role['role'] == 'member':
                can_edit = True

        if not can_edit:
            return jsonify({'success': False, 'error': 'У вас нет прав для редактирования этого поста'})

        conn.execute('UPDATE group_posts SET content = ? WHERE id = ?', (new_content, post_id))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/group_post/delete/<int:post_id>', methods=['POST'])
def delete_group_post(post_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    user_id = session['user_id']

    conn = get_db_connection()

    try:
        post = conn.execute('''
            SELECT gp.*, g.creator_id, gp.author_id as post_author_id
            FROM group_posts gp
            JOIN groups g ON gp.group_id = g.id
            WHERE gp.id = ?
        ''', (post_id,)).fetchone()

        if not post:
            return jsonify({'success': False, 'error': 'Post not found'})

        user_role = conn.execute('''
            SELECT role FROM group_members 
            WHERE group_id = ? AND user_id = ?
        ''', (post['group_id'], user_id)).fetchone()

        author_role = conn.execute('''
            SELECT role FROM group_members 
            WHERE group_id = ? AND user_id = ?
        ''', (post['group_id'], post['post_author_id'])).fetchone()

        can_delete = False

        if post['post_author_id'] == user_id:
            can_delete = True
        elif user_role and user_role['role'] == 'admin':
            can_delete = True
        elif user_role and user_role['role'] == 'moderator':
            if author_role and author_role['role'] == 'member':
                can_delete = True

        if not can_delete:
            return jsonify({'success': False, 'error': 'Not authorized to delete this post'})

        conn.execute('DELETE FROM group_posts WHERE id = ?', (post_id,))
        conn.execute('DELETE FROM group_post_likes WHERE post_id = ?', (post_id,))
        conn.execute('DELETE FROM post_media WHERE group_post_id = ?', (post_id,))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})

    finally:
        conn.close()


# Добавьте в main.py после импортов

@app.route('/friends/list')
def friends_list():
    """Получение списка друзей пользователя"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    conn = get_db_connection()

    friends = rows_to_dicts(conn.execute('''
        SELECT u.id, u.username, up.full_name, up.avatar
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.id IN (
            SELECT CASE 
                WHEN sender_id = ? THEN receiver_id 
                ELSE sender_id 
            END as friend_id
            FROM friendships 
            WHERE (sender_id = ? OR receiver_id = ?) AND status = 'accepted'
        )
        ORDER BY up.full_name, u.username
    ''', (user_id, user_id, user_id)).fetchall())

    conn.close()

    return jsonify({'success': True, 'friends': friends})


@app.route('/group/<int:group_id>/members')
def group_members_list(group_id):
    """Получение списка ID участников группы"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    conn = get_db_connection()
    members = conn.execute('''
        SELECT user_id FROM group_members WHERE group_id = ?
    ''', (group_id,)).fetchall()

    member_ids = [m['user_id'] for m in members]
    conn.close()

    return jsonify({'success': True, 'member_ids': member_ids})


@app.route('/group/<int:group_id>/invite', methods=['POST'])
def invite_to_group(group_id):
    """Приглашение друзей в группу"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    data = request.get_json()
    user_ids = data.get('user_ids', [])

    if not user_ids:
        return jsonify({'success': False, 'error': 'Нет пользователей для приглашения'})

    conn = get_db_connection()

    # Проверяем, является ли пользователь администратором группы
    role = conn.execute('''
        SELECT role FROM group_members 
        WHERE group_id = ? AND user_id = ?
    ''', (group_id, user_id)).fetchone()

    if not role or role['role'] != 'admin':
        conn.close()
        return jsonify({'success': False, 'error': 'Только администраторы могут приглашать'})

    sent_count = 0
    sent_ids = []

    for invited_id in user_ids:
        # Проверяем, не состоит ли уже пользователь в группе
        existing = conn.execute('''
            SELECT id FROM group_members 
            WHERE group_id = ? AND user_id = ?
        ''', (group_id, invited_id)).fetchone()

        if existing:
            continue

        # Проверяем, есть ли уже заявка
        existing_request = conn.execute('''
            SELECT id FROM group_requests 
            WHERE group_id = ? AND user_id = ? AND status = 'pending'
        ''', (group_id, invited_id)).fetchone()

        if existing_request:
            continue

        # Добавляем приглашение (создаем запись в group_requests)
        conn.execute('''
            INSERT INTO group_requests (group_id, user_id, status, created_at)
            VALUES (?, ?, 'pending', ?)
        ''', (group_id, invited_id, get_current_datetime()))

        sent_count += 1
        sent_ids.append(invited_id)

    conn.commit()
    conn.close()

    if sent_count > 0:
        return jsonify({
            'success': True,
            'message': f'Приглашения отправлены {sent_count} пользователям',
            'sent_count': sent_count,
            'sent_ids': sent_ids
        })
    else:
        return jsonify({'success': False, 'error': 'Все выбранные пользователи уже в группе или имеют активные заявки'})


@app.route('/feed', methods=['GET', 'POST'])
def feed():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']

    # Создание поста через форму в ленте
    if request.method == 'POST':
        content = request.form.get('content', '').strip()
        visibility = request.form.get('visibility', 'public')
        if content:
            conn = get_db_connection()
            current_datetime = get_current_datetime()
            conn.execute(
                'INSERT INTO posts (user_id, content, visibility, created_at) VALUES (?, ?, ?, ?)',
                (user_id, content, visibility, current_datetime)
            )
            conn.commit()
            conn.close()
            flash('Пост опубликован!', 'success')
        return redirect('/feed')

    conn = get_db_connection()

    user = conn.execute('SELECT username FROM users WHERE id = ?', (user_id,)).fetchone()
    username = user['username'] if user else ''

    try:
        friend_requests_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM friendships WHERE receiver_id = ? AND status = 'pending'",
            (user_id,)
        ).fetchone()['cnt']
    except Exception:
        friend_requests_count = 0

    my_groups = rows_to_dicts(conn.execute('''
        SELECT g.*, COUNT(DISTINCT gm.user_id) as members_count
        FROM groups g
        JOIN group_members gm ON g.id = gm.group_id
        WHERE gm.user_id = ?
        GROUP BY g.id
        ORDER BY g.name
        LIMIT 10
    ''', (user_id,)).fetchall())

    conn.close()

    feed_items = get_posts_feed(user_id=user_id, limit=50)

    return render_template('feed.html',
                           feed_items=feed_items,
                           username=username,
                           friend_requests_count=friend_requests_count,
                           my_groups=my_groups)


@app.route('/create_post', methods=['POST'])
def create_post():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    content = request.form.get('content', '').strip()
    visibility = request.form.get('visibility', 'public')

    if not content:
        flash('Пост не может быть пустым', 'error')
        return redirect('/feed')

    conn = get_db_connection()
    current_datetime = get_current_datetime()
    conn.execute('''
        INSERT INTO posts (user_id, content, visibility, created_at)
        VALUES (?, ?, ?, ?)
    ''', (user_id, content, visibility, current_datetime))
    conn.commit()
    conn.close()

    flash('Пост опубликован!', 'success')
    return redirect('/feed')


@app.route('/like_post/<int:post_id>', methods=['POST'])
def like_post(post_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    user_id = session['user_id']
    post_type = request.args.get('type', 'personal')

    conn = get_db_connection()

    try:
        if post_type == 'personal':
            existing_like = conn.execute('''
                SELECT id FROM post_likes WHERE post_id = ? AND user_id = ?
            ''', (post_id, user_id)).fetchone()

            if existing_like:
                conn.execute('DELETE FROM post_likes WHERE id = ?', (existing_like['id'],))
                action = 'unliked'
            else:
                conn.execute('INSERT INTO post_likes (post_id, user_id) VALUES (?, ?)', (post_id, user_id))
                action = 'liked'

            likes_count = \
                conn.execute('SELECT COUNT(*) as count FROM post_likes WHERE post_id = ?', (post_id,)).fetchone()[
                    'count']
        else:
            existing_like = conn.execute('''
                SELECT id FROM group_post_likes WHERE post_id = ? AND user_id = ?
            ''', (post_id, user_id)).fetchone()

            if existing_like:
                conn.execute('DELETE FROM group_post_likes WHERE id = ?', (existing_like['id'],))
                action = 'unliked'
            else:
                conn.execute('INSERT INTO group_post_likes (post_id, user_id) VALUES (?, ?)', (post_id, user_id))
                action = 'liked'

            likes_count = \
                conn.execute('SELECT COUNT(*) as count FROM group_post_likes WHERE post_id = ?', (post_id,)).fetchone()[
                    'count']

        conn.commit()
        return jsonify({'success': True, 'action': action, 'likes_count': likes_count})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/delete_post/<int:post_id>', methods=['POST'])
def delete_post(post_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authorized'})

    user_id = session['user_id']
    post_type = request.args.get('type', 'personal')

    conn = get_db_connection()

    try:
        if post_type == 'personal':
            post = conn.execute('SELECT user_id FROM posts WHERE id = ?', (post_id,)).fetchone()
            if not post or post['user_id'] != user_id:
                return jsonify({'success': False, 'error': 'Not authorized to delete this post'})

            conn.execute('DELETE FROM posts WHERE id = ?', (post_id,))
            conn.execute('DELETE FROM post_likes WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM comments WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM post_media WHERE post_id = ?', (post_id,))
        else:
            post = conn.execute('SELECT author_id, group_id FROM group_posts WHERE id = ?', (post_id,)).fetchone()
            if not post:
                return jsonify({'success': False, 'error': 'Post not found'})

            is_author = post['author_id'] == user_id
            is_admin = conn.execute('''
                SELECT role FROM group_members 
                WHERE group_id = ? AND user_id = ? AND role IN ('admin', 'moderator')
            ''', (post['group_id'], user_id)).fetchone()

            if not is_author and not is_admin:
                return jsonify({'success': False, 'error': 'Not authorized to delete this post'})

            conn.execute('DELETE FROM group_posts WHERE id = ?', (post_id,))
            conn.execute('DELETE FROM group_post_likes WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM post_media WHERE group_post_id = ?', (post_id,))

        conn.commit()
        return jsonify({'success': True})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/my_posts')
def my_posts():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем посты с количеством лайков и комментариев
    posts_rows = conn.execute('''
        SELECT p.*, 
               (SELECT COUNT(*) FROM post_likes WHERE post_id = p.id) as likes_count,
               (SELECT COUNT(*) FROM comments WHERE post_id = p.id) as comments_count
        FROM posts p 
        WHERE user_id = ? 
        ORDER BY created_at DESC
    ''', (user_id,)).fetchall()

    posts = rows_to_dicts(posts_rows)

    conn.close()

    return render_template('my_posts.html', posts=posts)


# ==================== ТЕХ-АДМИН ====================

@app.route('/techadmin')
def techadmin():
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь тех-админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()
    conn.close()

    if not user or user['role'] != 'techadmin':
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    return render_template('techadmin.html')


@app.route('/techadmin/reports')
def techadmin_reports():
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь тех-админом
    conn = get_db_connection()
    user = conn.execute('SELECT username, role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    print(f"=== ТЕХ-АДМИН ЗАПРОС ===")
    print(f"Пользователь: {user['username'] if user else 'не найден'}")
    print(f"Роль: {user['role'] if user else 'нет'}")

    if not user or user['role'] != 'techadmin':
        conn.close()
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    # Получаем статус фильтра
    status_filter = request.args.get('status', 'all')

    print(f"Фильтр статуса: {status_filter}")

    # Формируем запрос в зависимости от фильтра
    query = '''
        SELECT r.*, 
               reporter.username as reporter_username,
               reported.username as reported_username,
               reporter_profile.full_name as reporter_full_name,
               reported_profile.full_name as reported_full_name
        FROM reports r
        JOIN users reporter ON r.reporter_id = reporter.id
        JOIN users reported ON r.reported_id = reported.id
        LEFT JOIN user_profiles reporter_profile ON reporter.id = reporter_profile.user_id
        LEFT JOIN user_profiles reported_profile ON reported.id = reported_profile.user_id
    '''

    if status_filter != 'all':
        query += f" WHERE r.status = '{status_filter}'"

    query += " ORDER BY r.created_at DESC"

    print(f"Выполняем SQL запрос...")

    reports_rows = conn.execute(query).fetchall()
    reports = rows_to_dicts(reports_rows)

    print(f"Найдено жалоб: {len(reports)}")
    for i, report in enumerate(reports):
        print(
            f"  {i + 1}. ID: {report['id']}, Статус: {report['status']}, От: {report['reporter_username']}, На: {report['reported_username']}")

    conn.close()

    return render_template('techadmin_reports.html', reports=reports)


@app.route('/techadmin/report_action/<int:report_id>/<action>', methods=['POST'])
def techadmin_report_action(report_id, action):
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь тех-админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not user or user['role'] != 'techadmin':
        conn.close()
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    admin_notes = request.form.get('admin_notes', '')

    if action == 'approve':
        # Помечаем жалобу как обработанную
        conn.execute('''
            UPDATE reports 
            SET status = 'approved', admin_notes = ?, updated_at = CURRENT_TIMESTAMP
            WHERE id = ?
        ''', (admin_notes, report_id))

        # Получаем информацию о жалобе
        report = conn.execute('SELECT reported_id FROM reports WHERE id = ?', (report_id,)).fetchone()
        if report:
            # Баним пользователя (можно добавить более сложную логику)
            conn.execute('''
                UPDATE users SET is_banned = 1 WHERE id = ?
            ''', (report['reported_id'],))

        flash('Жалоба одобрена, пользователь забанен', 'success')


    elif action == 'reject':
        # Помечаем жалобу как отклоненную
        conn.execute('''
            UPDATE reports 
            SET status = 'rejected', admin_notes = ?, updated_at = CURRENT_TIMESTAMP
            WHERE id = ?
        ''', (admin_notes, report_id))
        flash('Жалоба отклонена', 'info')

    elif action == 'delete':
        # Удаляем жалобу
        conn.execute('DELETE FROM reports WHERE id = ?', (report_id,))
        flash('Жалоба удалена', 'info')

    conn.commit()
    conn.close()

    return redirect('/techadmin/reports')


@app.route('/techadmin/stats')
def techadmin_stats():
    """Получение статистики для тех-админа"""
    if 'user_id' not in session:
        return jsonify({'error': 'Требуется авторизация'}), 401

    conn = get_db_connection()

    # Проверяем, является ли пользователь тех-админом
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not user or user['role'] != 'techadmin':
        conn.close()
        return jsonify({'error': 'Доступ запрещен'}), 403

    # Получаем статистику жалоб
    pending_reports = conn.execute('''SELECT COUNT(*) as count FROM reports WHERE status = 'pending''').fetchone()[
        'count']
    total_reports = conn.execute('''SELECT COUNT(*) as count FROM reports''').fetchone()['count']

    # Получаем статистику пользователей
    banned_users = conn.execute('''SELECT COUNT(*) as count FROM users WHERE is_banned = 1''').fetchone()['count']
    total_users = conn.execute('''SELECT COUNT(*) as count FROM users''').fetchone()['count']

    conn.close()

    return jsonify({
        'pending_reports': pending_reports,
        'total_reports': total_reports,
        'banned_users': banned_users,
        'total_users': total_users
    })


# ==================== АУТЕНТИФИКАЦИЯ ====================

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        agree_terms = request.form.get('agree_terms')
        if not agree_terms:
            return render_template('register.html', error="Вы должны принять услвоия пользовательского соглашения")
        username = request.form['username']
        password = request.form['password']
        confirm_password = request.form.get('confirm_password', '')

        if password != confirm_password:
            return render_template('register.html', error="Пароли не совпадают")

        if len(password) < 6:
            return render_template('register.html', error="Пароль должен содержать минимум 6 символов")

        try:
            connection = get_db_connection()
            cursor = connection.cursor()
            cursor.execute('INSERT INTO users (username, password) VALUES (?, ?)', (username, password))
            connection.commit()
            connection.close()
            flash('Регистрация успешна! Теперь вы можете войти.', 'success')
            return redirect('/login')
        except sqlite3.IntegrityError:
            return render_template('register.html', error="Пользователь с таким именем уже существует")
        except Exception as e:
            print(f"Ошибка регистрации: {e}")
            return render_template('register.html', error="Ошибка при регистрации")

    return render_template('register.html')


@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']

        try:
            connection = get_db_connection()
            cursor = connection.cursor()
            cursor.execute("SELECT id, username, role, is_banned FROM users WHERE username = ? AND password = ?",
                           (username, password))
            user_row = cursor.fetchone()

            if user_row:
                user = dict(user_row)

                # Проверяем, не забанен ли пользователь
                if user.get('role') == 'user' and user.get('is_banned', 0) == 1:
                    connection.close()
                    return render_template('login.html', error="Ваш аккаунт заблокирован")

                # Проверяем включена ли 2FA
                tfa = connection.execute('''
                    SELECT * FROM two_factor_auth
                    WHERE user_id = ? AND is_enabled = 1 AND telegram_chat_id IS NOT NULL
                ''', (user['id'],)).fetchone()

                if tfa:
                    tfa = dict(tfa)
                    # Генерируем одноразовый код (5 минут)
                    code = generate_code()
                    expires = (datetime.now() + timedelta(minutes=5)).strftime('%Y-%m-%d %H:%M:%S')
                    connection.execute('''
                        UPDATE two_factor_auth SET auth_code = ?, auth_code_expires = ?
                        WHERE user_id = ?
                    ''', (code, expires, user['id']))
                    connection.commit()
                    connection.close()

                    # Отправляем код в Telegram
                    send_telegram_2fa_code(tfa['telegram_chat_id'], code, user['username'])

                    # Временная сессия — только для 2FA шага
                    session['2fa_user_id'] = user['id']
                    return redirect('/2fa/verify')
                else:
                    connection.close()
                    session['user_id'] = user['id']
                    session['username'] = user['username']
                    return redirect_based_on_role(user['username'])
            else:
                connection.close()
                return render_template('login.html', error="Неверное имя пользователя или пароль")
        except Exception as e:
            print(f"Ошибка входа: {e}")
            return render_template('login.html', error="Ошибка базы данных")

    return render_template('login.html')


@app.route('/terms')
def terms():
    return render_template('terms.html')


# ==================== НАСТРОЙКИ АККАУНТА ====================

@app.route('/settings')
def settings():
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()
    user = row_to_dict(conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone())
    tfa = row_to_dict(conn.execute('SELECT * FROM two_factor_auth WHERE user_id = ?', (user_id,)).fetchone())
    conn.close()

    return render_template('settings.html', user=user, tfa=tfa)


@app.route('/settings/change_password', methods=['POST'])
def settings_change_password():
    """Смена пароля из настроек"""
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    current_password = request.form.get('current_password', '')
    new_password = request.form.get('new_password', '')
    confirm_password = request.form.get('confirm_password', '')

    conn = get_db_connection()
    user = row_to_dict(conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone())

    if user['password'] != current_password:
        conn.close()
        flash('Неверный текущий пароль', 'error')
        return redirect('/settings')

    if new_password != confirm_password:
        conn.close()
        flash('Новые пароли не совпадают', 'error')
        return redirect('/settings')

    if len(new_password) < 6:
        conn.close()
        flash('Пароль должен содержать минимум 6 символов', 'error')
        return redirect('/settings')

    conn.execute('UPDATE users SET password = ? WHERE id = ?', (new_password, user_id))
    conn.commit()
    conn.close()
    flash('Пароль успешно изменён', 'success')
    return redirect('/settings')


# ==================== 2FA: НАЧАЛО ПРИВЯЗКИ ====================

@app.route('/settings/2fa/start', methods=['POST'])
def settings_2fa_start():
    """Генерирует код привязки и перенаправляет на инструкцию"""
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()

    # Если уже привязан — ничего не делаем
    tfa = row_to_dict(conn.execute('SELECT * FROM two_factor_auth WHERE user_id = ?', (user_id,)).fetchone())
    if tfa and tfa.get('is_enabled') and tfa.get('telegram_chat_id'):
        conn.close()
        flash('Telegram уже привязан к вашему аккаунту', 'info')
        return redirect('/settings')

    # Генерируем код привязки на 15 минут
    link_code = generate_link_code()
    expires = (datetime.now() + timedelta(minutes=15)).strftime('%Y-%m-%d %H:%M:%S')

    if tfa:
        conn.execute('''
            UPDATE two_factor_auth
            SET link_code = ?, link_code_expires = ?, telegram_chat_id = NULL, is_enabled = 0
            WHERE user_id = ?
        ''', (link_code, expires, user_id))
    else:
        conn.execute('''
            INSERT INTO two_factor_auth (user_id, link_code, link_code_expires)
            VALUES (?, ?, ?)
        ''', (user_id, link_code, expires))

    conn.commit()
    conn.close()

    bot_link = f"https://t.me/{BOT_USERNAME}?start={link_code}"
    return render_template('2fa_setup.html', bot_link=bot_link, link_code=link_code)


@app.route('/settings/2fa/check_status')
def settings_2fa_check_status():
    """AJAX: проверяет, привязал ли пользователь Telegram"""
    if 'user_id' not in session:
        return jsonify({'enabled': False})

    user_id = session['user_id']
    conn = get_db_connection()
    tfa = row_to_dict(conn.execute(
        'SELECT is_enabled, telegram_chat_id FROM two_factor_auth WHERE user_id = ?', (user_id,)
    ).fetchone())
    conn.close()

    enabled = bool(tfa and tfa.get('is_enabled') and tfa.get('telegram_chat_id'))
    return jsonify({'enabled': enabled})


# ==================== 2FA: ОТКЛЮЧЕНИЕ ====================

@app.route('/settings/2fa/disable', methods=['POST'])
def settings_2fa_disable():
    """Отключает 2FA и удаляет привязку Telegram"""
    if 'user_id' not in session:
        return redirect('/login')

    user_id = session['user_id']
    conn = get_db_connection()
    conn.execute('''
        UPDATE two_factor_auth
        SET is_enabled = 0, telegram_chat_id = NULL,
            link_code = NULL, link_code_expires = NULL
        WHERE user_id = ?
    ''', (user_id,))
    conn.commit()
    conn.close()

    flash('Двухфакторная аутентификация отключена', 'success')
    return redirect('/settings')


# ==================== 2FA: ПРОВЕРКА КОДА ПРИ ВХОДЕ ====================

@app.route('/2fa/verify', methods=['GET', 'POST'])
def twofa_verify():
    """Страница ввода кода после успешного ввода пароля"""
    if '2fa_user_id' not in session:
        return redirect('/login')

    user_id = session['2fa_user_id']

    if request.method == 'POST':
        entered_code = request.form.get('code', '').strip()

        conn = get_db_connection()
        tfa = conn.execute('''
            SELECT * FROM two_factor_auth
            WHERE user_id = ?
              AND auth_code = ?
              AND auth_code_expires > datetime('now')
        ''', (user_id, entered_code)).fetchone()

        if tfa:
            # Сбрасываем одноразовый код
            conn.execute('''
                UPDATE two_factor_auth SET auth_code = NULL, auth_code_expires = NULL
                WHERE user_id = ?
            ''', (user_id,))
            conn.commit()

            user = row_to_dict(conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone())
            conn.close()

            # Полноценный вход
            session.pop('2fa_user_id', None)
            session['user_id'] = user['id']
            session['username'] = user['username']
            return redirect_based_on_role(user['username'])
        else:
            conn.close()
            return render_template('2fa_verify.html', error='Неверный или истёкший код')

    return render_template('2fa_verify.html')


@app.route('/2fa/resend', methods=['POST'])
def twofa_resend():
    """Повторная отправка кода 2FA"""
    if '2fa_user_id' not in session:
        return jsonify({'success': False, 'error': 'Сессия истекла'}), 401

    user_id = session['2fa_user_id']
    conn = get_db_connection()

    tfa = row_to_dict(conn.execute(
        'SELECT * FROM two_factor_auth WHERE user_id = ? AND is_enabled = 1', (user_id,)
    ).fetchone())

    if not tfa or not tfa.get('telegram_chat_id'):
        conn.close()
        return jsonify({'success': False, 'error': 'Telegram не привязан'}), 400

    user = row_to_dict(conn.execute('SELECT username FROM users WHERE id = ?', (user_id,)).fetchone())

    code = generate_code()
    expires = (datetime.now() + timedelta(minutes=5)).strftime('%Y-%m-%d %H:%M:%S')
    conn.execute('''
        UPDATE two_factor_auth SET auth_code = ?, auth_code_expires = ?
        WHERE user_id = ?
    ''', (code, expires, user_id))
    conn.commit()
    conn.close()

    ok = send_telegram_2fa_code(tfa['telegram_chat_id'], code, user['username'])
    if ok:
        return jsonify({'success': True})
    else:
        return jsonify({'success': False, 'error': 'Не удалось отправить сообщение в Telegram'})


# ==================== TELEGRAM БОТ (запускается в фоновом потоке) ====================

def run_telegram_bot():
    """Запускает Telegram-бот в отдельном потоке с собственным event loop"""
    try:
        from telegram import Update
        from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes
    except ImportError:
        print("[BOT] python-telegram-bot не установлен. Установите: pip install python-telegram-bot")
        return

    if BOT_TOKEN == "YOUR_BOT_TOKEN":
        print("[BOT] Токен не задан — бот не запущен. Укажи BOT_TOKEN в main.py")
        return

    async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
        chat_id = update.effective_chat.id
        args = context.args

        if args:
            link_code = args[0]
            conn = get_db_connection()
            try:
                record = conn.execute('''
                    SELECT user_id FROM two_factor_auth
                    WHERE link_code = ?
                      AND link_code_expires > datetime('now')
                      AND telegram_chat_id IS NULL
                ''', (link_code,)).fetchone()

                if record:
                    user_id = record['user_id']
                    user = conn.execute('SELECT username FROM users WHERE id = ?', (user_id,)).fetchone()
                    username = user['username'] if user else f"id{user_id}"

                    conn.execute('''
                        UPDATE two_factor_auth
                        SET telegram_chat_id = ?,
                            link_code = NULL,
                            link_code_expires = NULL,
                            is_enabled = 1
                        WHERE user_id = ?
                    ''', (str(chat_id), user_id))
                    conn.commit()

                    await update.message.reply_text(
                        f"✅ *Telegram успешно привязан к аккаунту {username}!*\n\n"
                        f"Двухфакторная аутентификация активирована.\n"
                        f"При каждом входе я буду отправлять тебе код подтверждения.",
                        parse_mode='Markdown'
                    )
                else:
                    await update.message.reply_text(
                        "❌ *Код привязки недействителен или истёк.*\n\n"
                        "Вернись в настройки аккаунта и запроси новый код.",
                        parse_mode='Markdown'
                    )
            except Exception as e:
                print(f"[BOT] Ошибка привязки: {e}")
                await update.message.reply_text("⚠️ Произошла ошибка. Попробуй позже.")
            finally:
                conn.close()
        else:
            await update.message.reply_text(
                "👋 *Привет! Я бот для двухфакторной аутентификации.*\n\n"
                "Чтобы привязать аккаунт, зайди в *Настройки* на сайте "
                "и нажми «Подключить Telegram».",
                parse_mode='Markdown'
            )

    async def status_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
        chat_id = str(update.effective_chat.id)
        conn = get_db_connection()
        try:
            record = conn.execute('''
                SELECT u.username FROM two_factor_auth tfa
                JOIN users u ON tfa.user_id = u.id
                WHERE tfa.telegram_chat_id = ? AND tfa.is_enabled = 1
            ''', (chat_id,)).fetchone()

            if record:
                await update.message.reply_text(
                    f"✅ Этот чат привязан к аккаунту *{record['username']}*.\n"
                    f"Двухфакторная аутентификация активна.",
                    parse_mode='Markdown'
                )
            else:
                await update.message.reply_text(
                    "ℹ️ Этот чат не привязан ни к одному аккаунту.",
                    parse_mode='Markdown'
                )
        finally:
            conn.close()

    async def help_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
        await update.message.reply_text(
            "📋 *Команды:*\n/start — привязка аккаунта\n/status — проверить привязку\n/help — помощь",
            parse_mode='Markdown'
        )

    async def unknown(update: Update, context: ContextTypes.DEFAULT_TYPE):
        await update.message.reply_text("ℹ️ Используй /help для справки.")

    # Создаём собственный event loop для потока
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)

    async def main_bot():
        bot_app = ApplicationBuilder().token(BOT_TOKEN).build()
        bot_app.add_handler(CommandHandler("start", start))
        bot_app.add_handler(CommandHandler("status", status_cmd))
        bot_app.add_handler(CommandHandler("help", help_cmd))
        bot_app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, unknown))

        print("[BOT] Telegram-бот запущен в фоновом потоке")
        async with bot_app:
            await bot_app.initialize()
            await bot_app.start()
            await bot_app.updater.start_polling(drop_pending_updates=True)
            while True:
                await asyncio.sleep(1)

    loop.run_until_complete(main_bot())


# ==================== АДМИН-ПАНЕЛЬ ====================

@app.route('/admin')
def admin_panel():
    """Главная страница админ-панели"""
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()
    conn.close()

    if not user or user['role'] != 'admin':
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    return redirect('/admin/users')  # По умолчанию переходим к управлению пользователями


@app.route('/admin/users')
def admin_users():
    """Управление пользователями"""
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not user or user['role'] != 'admin':
        conn.close()
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    # Получаем параметры поиска и фильтрации
    search_query = request.args.get('search', '').strip()
    role_filter = request.args.get('role', 'all')
    banned_filter = request.args.get('banned', 'all')

    # Формируем запрос
    query = '''
        SELECT u.*, 
               up.full_name, 
               up.avatar,
               (SELECT COUNT(*) FROM posts WHERE user_id = u.id) as posts_count,
               (SELECT COUNT(*) FROM friendships WHERE (sender_id = u.id OR receiver_id = u.id) AND status = 'accepted') as friends_count
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE 1=1
    '''

    params = []

    # Применяем фильтры
    if search_query:
        query += ' AND (u.username LIKE ? OR up.full_name LIKE ?)'
        params.extend([f'%{search_query}%', f'%{search_query}%'])

    if role_filter != 'all':
        query += ' AND u.role = ?'
        params.append(role_filter)

    if banned_filter != 'all':
        if banned_filter == 'banned':
            query += ' AND u.is_banned = 1'
        elif banned_filter == 'active':
            query += ' AND u.is_banned = 0'

    query += ' ORDER BY u.created_at DESC'

    users = rows_to_dicts(conn.execute(query, params).fetchall())

    # Получаем статистику
    total_users = conn.execute('SELECT COUNT(*) as count FROM users').fetchone()['count']
    admin_users = conn.execute("SELECT COUNT(*) as count FROM users WHERE role = 'admin'").fetchone()['count']
    techadmin_users = conn.execute("SELECT COUNT(*) as count FROM users WHERE role = 'techadmin'").fetchone()['count']
    banned_users = conn.execute("SELECT COUNT(*) as count FROM users WHERE is_banned = 1").fetchone()['count']

    conn.close()

    return render_template('admin_users.html',
                           users=users,
                           search_query=search_query,
                           role_filter=role_filter,
                           banned_filter=banned_filter,
                           total_users=total_users,
                           admin_users=admin_users,
                           techadmin_users=techadmin_users,
                           banned_users=banned_users)


@app.route('/admin/banned')
def admin_banned():
    """Забаненные пользователи"""
    if 'user_id' not in session:
        return redirect('/login')

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not user or user['role'] != 'admin':
        conn.close()
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    # Получаем список забаненных пользователей
    users = rows_to_dicts(conn.execute('''
        SELECT u.*, 
               up.full_name, 
               up.avatar,
               (SELECT COUNT(*) FROM posts WHERE user_id = u.id) as posts_count
        FROM users u
        LEFT JOIN user_profiles up ON u.id = up.user_id
        WHERE u.is_banned = 1
        ORDER BY u.created_at DESC
    ''').fetchall())

    conn.close()

    return render_template('admin_banned.html', users=users)


@app.route('/admin/change_role/<int:user_id>', methods=['POST'])
def admin_change_role(user_id):
    """Изменение роли пользователя"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    admin_user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not admin_user or admin_user['role'] != 'admin':
        conn.close()
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    data = request.get_json()
    new_role = data.get('role')

    if not new_role or new_role not in ['admin', 'techadmin', 'user']:
        return jsonify({'success': False, 'error': 'Недопустимая роль'})

    # Получаем информацию о пользователе
    target_user = conn.execute('SELECT id, username, role FROM users WHERE id = ?', (user_id,)).fetchone()

    if not target_user:
        return jsonify({'success': False, 'error': 'Пользователь не найден'})

    # Не позволяем менять роль главного админа
    if target_user['username'] == 'admin':
        return jsonify({'success': False, 'error': 'Нельзя изменить роль главного администратора'})

    # Не позволяем менять свою собственную роль
    if user_id == session['user_id']:
        return jsonify({'success': False, 'error': 'Нельзя изменить свою собственную роль'})

    old_role = target_user['role']

    # Изменяем роль
    conn.execute('UPDATE users SET role = ? WHERE id = ?', (new_role, user_id))
    conn.commit()
    conn.close()

    # Определяем название роли для сообщения
    role_names = {
        'admin': 'администратора',
        'techadmin': 'технического администратора',
        'user': 'обычного пользователя'
    }

    role_name = role_names.get(new_role, new_role)
    old_role_name = role_names.get(old_role, old_role)

    return jsonify({
        'success': True,
        'message': f'Роль пользователя изменена с {old_role_name} на {role_name}',
        'new_role': new_role,
        'old_role': old_role
    })


@app.route('/admin/ban_user/<int:user_id>', methods=['POST'])
def admin_ban_user(user_id):
    """Бан пользователя"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    admin_user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not admin_user or admin_user['role'] != 'admin':
        conn.close()
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    data = request.get_json()
    ban_reason = data.get('reason', '')

    # Получаем информацию о пользователе
    target_user = conn.execute('SELECT id, username, role FROM users WHERE id = ?', (user_id,)).fetchone()

    if not target_user:
        return jsonify({'success': False, 'error': 'Пользователь не найден'})

    # Не позволяем забанить другого админа или тех-админа
    if target_user['role'] in ['admin', 'techadmin']:
        return jsonify({'success': False, 'error': 'Нельзя забанить администратора или тех-админа'})

    # Не позволяем забанить себя
    if user_id == session['user_id']:
        return jsonify({'success': False, 'error': 'Нельзя забанить себя'})

    # Баним пользователя
    conn.execute('UPDATE users SET is_banned = 1 WHERE id = ?', (user_id,))

    # Добавляем запись в журнал банов (можно создать таблицу ban_history если нужно)
    conn.commit()
    conn.close()

    return jsonify({'success': True, 'message': f'Пользователь {target_user["username"]} забанен'})


@app.route('/admin/unban_user/<int:user_id>', methods=['POST'])
def admin_unban_user(user_id):
    """Разбан пользователя"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    admin_user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not admin_user or admin_user['role'] != 'admin':
        conn.close()
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    # Получаем информацию о пользователе
    target_user = conn.execute('SELECT id, username FROM users WHERE id = ?', (user_id,)).fetchone()

    if not target_user:
        return jsonify({'success': False, 'error': 'Пользователь не найден'})

    # Разбаниваем пользователя
    conn.execute('UPDATE users SET is_banned = 0 WHERE id = ?', (user_id,))
    conn.commit()
    conn.close()

    return jsonify({'success': True, 'message': f'Пользователь {target_user["username"]} разбанен'})


@app.route('/admin/get_user_stats')
def admin_get_user_stats():
    """Получение статистики пользователей для админа"""
    if 'user_id' not in session:
        return jsonify({'error': 'Требуется авторизация'}), 401

    # Проверяем, является ли пользователь админом
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (session['user_id'],)).fetchone()

    if not user or user['role'] != 'admin':
        conn.close()
        return jsonify({'error': 'Доступ запрещен'}), 403

    # Получаем статистику
    today = datetime.now().strftime('%Y-%m-%d')

    total_users = conn.execute('SELECT COUNT(*) as count FROM users').fetchone()['count']
    today_users = conn.execute('SELECT COUNT(*) as count FROM users WHERE DATE(created_at) = ?', (today,)).fetchone()[
        'count']
    banned_users = conn.execute("SELECT COUNT(*) as count FROM users WHERE is_banned = 1").fetchone()['count']
    active_users = total_users - banned_users

    # Получаем распределение по ролям
    roles_stats = rows_to_dicts(conn.execute('''
        SELECT role, COUNT(*) as count 
        FROM users 
        GROUP BY role
    ''').fetchall())

    conn.close()

    return jsonify({
        'total_users': total_users,
        'today_users': today_users,
        'banned_users': banned_users,
        'active_users': active_users,
        'roles_stats': roles_stats
    })


@app.route('/admin/groups')
def admin_groups():
    """Управление группами"""
    if 'user_id' not in session:
        return redirect('/login')
    if not is_admin(session['user_id']):
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    search_query = request.args.get('search', '').strip()
    conn = get_db_connection()

    query = """
        SELECT g.*,
               u.username as creator_username,
               COUNT(DISTINCT gm.user_id) as members_count,
               COUNT(DISTINCT gp.id)      as posts_count
        FROM groups g
        LEFT JOIN users u ON g.creator_id = u.id
        LEFT JOIN group_members gm ON gm.group_id = g.id
        LEFT JOIN group_posts gp ON gp.group_id = g.id
    """
    params = []
    if search_query:
        query += " WHERE g.name LIKE ? OR g.description LIKE ?"
        params.extend([f'%{search_query}%', f'%{search_query}%'])
    query += " GROUP BY g.id ORDER BY members_count DESC"

    groups = rows_to_dicts(conn.execute(query, params).fetchall())
    total_groups  = conn.execute('SELECT COUNT(*) as c FROM groups').fetchone()['c']
    total_members = conn.execute('SELECT COUNT(*) as c FROM group_members').fetchone()['c']
    total_posts   = conn.execute('SELECT COUNT(*) as c FROM group_posts').fetchone()['c']
    conn.close()

    return render_template('admin_groups.html',
                           groups=groups,
                           search_query=search_query,
                           total_groups=total_groups,
                           total_members=total_members,
                           total_posts=total_posts)


@app.route('/admin/posts')
def admin_posts():
    """Управление постами"""
    if 'user_id' not in session:
        return redirect('/login')
    if not is_admin(session['user_id']):
        flash('Доступ запрещен', 'error')
        return redirect('/home')

    search_query = request.args.get('search', '').strip()
    post_type    = request.args.get('type', 'all')   # all | personal | group
    conn = get_db_connection()

    personal_posts = []
    group_posts    = []

    if post_type in ('all', 'personal'):
        q = """
            SELECT p.id, p.content, p.created_at, p.user_id,
                   u.username as author_username,
                   COALESCE(up.full_name, u.username) as author_name,
                   (SELECT COUNT(*) FROM post_likes  WHERE post_id = p.id) as likes_count,
                   (SELECT COUNT(*) FROM comments    WHERE post_id = p.id) as comments_count
            FROM posts p
            JOIN users u ON p.user_id = u.id
            LEFT JOIN user_profiles up ON up.user_id = u.id
        """
        p = []
        if search_query:
            q += " WHERE p.content LIKE ? OR u.username LIKE ?"
            p.extend([f'%{search_query}%', f'%{search_query}%'])
        q += " ORDER BY p.created_at DESC LIMIT 200"
        personal_posts = rows_to_dicts(conn.execute(q, p).fetchall())

    if post_type in ('all', 'group'):
        q = """
            SELECT gp.id, gp.content, gp.created_at, gp.group_id, gp.author_id,
                   u.username as author_username,
                   COALESCE(up.full_name, u.username) as author_name,
                   g.name as group_name,
                   (SELECT COUNT(*) FROM group_post_likes WHERE post_id = gp.id) as likes_count
            FROM group_posts gp
            JOIN users u  ON gp.author_id = u.id
            LEFT JOIN user_profiles up ON up.user_id = u.id
            JOIN groups g ON gp.group_id = g.id
        """
        p = []
        if search_query:
            q += " WHERE gp.content LIKE ? OR u.username LIKE ? OR g.name LIKE ?"
            p.extend([f'%{search_query}%', f'%{search_query}%', f'%{search_query}%'])
        q += " ORDER BY gp.created_at DESC LIMIT 200"
        group_posts = rows_to_dicts(conn.execute(q, p).fetchall())

    total_personal = conn.execute('SELECT COUNT(*) as c FROM posts').fetchone()['c']
    total_group    = conn.execute('SELECT COUNT(*) as c FROM group_posts').fetchone()['c']
    conn.close()

    return render_template('admin_posts.html',
                           personal_posts=personal_posts,
                           group_posts=group_posts,
                           search_query=search_query,
                           post_type=post_type,
                           total_personal=total_personal,
                           total_group=total_group)


def is_admin(user_id):
    """Проверяет, является ли пользователь администратором"""
    conn = get_db_connection()
    user = conn.execute('SELECT role FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    return user and user['role'] == 'admin'


@app.route('/admin/delete_group/<int:group_id>', methods=['POST'])
def admin_delete_group(group_id):
    """Удаление любой группы администратором"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    if not is_admin(session['user_id']):
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    conn = get_db_connection()
    try:
        group = conn.execute('SELECT name, avatar FROM groups WHERE id = ?', (group_id,)).fetchone()
        if not group:
            return jsonify({'success': False, 'error': 'Группа не найдена'})

        # Удаляем медиафайлы постов группы
        media_files = conn.execute('''
            SELECT filename FROM post_media
            WHERE group_post_id IN (SELECT id FROM group_posts WHERE group_id = ?)
        ''', (group_id,)).fetchall()
        for m in media_files:
            try:
                os.remove(os.path.join(app.config['POST_MEDIA_FOLDER'], m['filename']))
            except:
                pass

        # Удаляем аватар группы
        if group['avatar'] and group['avatar'] != 'default_group.png':
            try:
                os.remove(os.path.join(app.config['GROUP_UPLOAD_FOLDER'], group['avatar']))
            except:
                pass

        # Удаляем все связанные данные
        conn.execute('DELETE FROM post_media WHERE group_post_id IN (SELECT id FROM group_posts WHERE group_id = ?)', (group_id,))
        conn.execute('DELETE FROM group_post_likes WHERE post_id IN (SELECT id FROM group_posts WHERE group_id = ?)', (group_id,))
        conn.execute('DELETE FROM group_posts WHERE group_id = ?', (group_id,))
        conn.execute('DELETE FROM group_members WHERE group_id = ?', (group_id,))
        conn.execute('DELETE FROM group_requests WHERE group_id = ?', (group_id,))
        conn.execute('DELETE FROM groups WHERE id = ?', (group_id,))
        conn.commit()

        return jsonify({'success': True, 'message': f'Группа "{group["name"]}" удалена'})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/admin/delete_post/<int:post_id>', methods=['POST'])
def admin_delete_post(post_id):
    """Удаление любого поста администратором"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    if not is_admin(session['user_id']):
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    data = request.get_json() or {}
    is_group_post = data.get('is_group_post', False)

    conn = get_db_connection()
    try:
        if is_group_post:
            post = conn.execute('SELECT id FROM group_posts WHERE id = ?', (post_id,)).fetchone()
            if not post:
                return jsonify({'success': False, 'error': 'Пост не найден'})

            # Удаляем медиафайлы
            media_files = conn.execute('SELECT filename FROM post_media WHERE group_post_id = ?', (post_id,)).fetchall()
            for m in media_files:
                try:
                    os.remove(os.path.join(app.config['POST_MEDIA_FOLDER'], m['filename']))
                except:
                    pass

            conn.execute('DELETE FROM post_media WHERE group_post_id = ?', (post_id,))
            conn.execute('DELETE FROM group_post_likes WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM group_posts WHERE id = ?', (post_id,))
        else:
            post = conn.execute('SELECT id FROM posts WHERE id = ?', (post_id,)).fetchone()
            if not post:
                return jsonify({'success': False, 'error': 'Пост не найден'})

            # Удаляем медиафайлы
            media_files = conn.execute('SELECT filename FROM post_media WHERE post_id = ?', (post_id,)).fetchall()
            for m in media_files:
                try:
                    os.remove(os.path.join(app.config['POST_MEDIA_FOLDER'], m['filename']))
                except:
                    pass

            conn.execute('DELETE FROM post_media WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM post_likes WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM comments WHERE post_id = ?', (post_id,))
            conn.execute('DELETE FROM posts WHERE id = ?', (post_id,))

        conn.commit()
        return jsonify({'success': True, 'message': 'Пост удален'})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


@app.route('/admin/delete_comment/<int:comment_id>', methods=['POST'])
def admin_delete_comment(comment_id):
    """Удаление любого комментария администратором"""
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Требуется авторизация'}), 401

    if not is_admin(session['user_id']):
        return jsonify({'success': False, 'error': 'Доступ запрещен'}), 403

    conn = get_db_connection()
    try:
        comment = conn.execute('SELECT id, post_id FROM comments WHERE id = ?', (comment_id,)).fetchone()
        if not comment:
            return jsonify({'success': False, 'error': 'Комментарий не найден'})

        post_id = comment['post_id']
        conn.execute('DELETE FROM comments WHERE id = ?', (comment_id,))
        comments_count = conn.execute('SELECT COUNT(*) as count FROM comments WHERE post_id = ?', (post_id,)).fetchone()['count']
        conn.commit()

        return jsonify({'success': True, 'message': 'Комментарий удален', 'comments_count': comments_count})

    except Exception as e:
        conn.rollback()
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()


def redirect_based_on_role(username):
    """Перенаправляет пользователя на соответствующую страницу в зависимости от роли"""
    if username == 'admin':
        return redirect('/admin')
    elif username == 'techadmin':
        return redirect('/techadmin')
    else:
        # Проверяем роль в базе данных
        conn = get_db_connection()
        user = conn.execute('SELECT role FROM users WHERE username = ?', (username,)).fetchone()
        conn.close()

        if user:
            if user['role'] == 'admin':
                return redirect('/admin')
            elif user['role'] == 'techadmin':
                return redirect('/techadmin')

    return redirect('/home')


@app.route('/user')
def user_page():
    return redirect('/home')


@app.route('/logout')
def logout():
    session.clear()
    return redirect('/login')


@app.route('/get_user_stats')
def get_user_stats():
    if 'user_id' not in session:
        return jsonify({'error': 'Требуется авторизация'}), 401

    user_id = session['user_id']
    conn = get_db_connection()

    # Получаем количество постов пользователя
    posts_count = conn.execute('''
        SELECT COUNT(*) as count FROM posts WHERE user_id = ?
    ''', (user_id,)).fetchone()['count']

    # Получаем количество друзей
    friends_count = conn.execute('''
        SELECT COUNT(*) as count FROM friendships 
        WHERE (sender_id = ? OR receiver_id = ?) AND status = 'accepted'
    ''', (user_id, user_id)).fetchone()['count']

    conn.close()

    return jsonify({
        'posts_count': posts_count,
        'friends_count': friends_count
    })


@app.route('/feed/more')
def feed_more():
    """Подгрузка следующей порции постов (бесконечная лента)"""
    if 'user_id' not in session:
        return jsonify({'posts': [], 'has_more': False})

    user_id     = session['user_id']
    offset      = request.args.get('offset', 0, type=int)
    filter_type = request.args.get('filter', 'all')
    per_page    = 10

    try:
        items    = get_posts_feed(user_id, limit=per_page + 1, filter_type=filter_type, offset=offset)
        has_more = len(items) > per_page
        items    = items[:per_page]

        result = []
        for item in items:
            media = []
            for m in item.get('media_files', []):
                media.append({
                    'filename':          m.get('filename', ''),
                    'file_type':         m.get('file_type', 'image'),
                    'original_filename': m.get('original_filename', ''),
                    'id':                m.get('id', 0),
                })
            result.append({
                'id':             item['id'],
                'user_id':        item['user_id'],
                'username':       item.get('username', ''),
                'author_name':    item.get('author_name', item.get('username', '')),
                'author_avatar':  item.get('author_avatar', ''),
                'content':        item.get('content', ''),
                'created_at':     item.get('created_at', ''),
                'likes_count':    item.get('likes_count', 0),
                'comments_count': item.get('comments_count', 0),
                'is_liked':       bool(item.get('is_liked', 0)),
                'is_own':         item['user_id'] == user_id,
                'media':          media,
            })

        return jsonify({'posts': result, 'has_more': has_more})

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'posts': [], 'has_more': False, 'error': str(e)})

@app.route('/debug_video')
def debug_video():
    if 'user_id' not in session:
        return redirect('/login')

    conn = get_db_connection()

    # Получаем все видео из БД
    videos = conn.execute('''
        SELECT pm.*, p.content, p.user_id, u.username
        FROM post_media pm
        JOIN posts p ON pm.post_id = p.id
        JOIN users u ON p.user_id = u.id
        WHERE pm.file_type = 'video'
        ORDER BY pm.id DESC
    ''').fetchall()

    # Проверяем физическое существование файлов
    import os
    video_files = []
    missing_files = []

    for video in videos:
        file_path = os.path.join(app.config['POST_MEDIA_FOLDER'], video['filename'])
        if os.path.exists(file_path):
            file_size = os.path.getsize(file_path)
            video_files.append({
                'id': video['id'],
                'filename': video['filename'],
                'post_id': video['post_id'],
                'username': video['username'],
                'size': file_size,
                'exists': True
            })
        else:
            missing_files.append({
                'id': video['id'],
                'filename': video['filename'],
                'post_id': video['post_id']
            })

    # Получаем все физические видеофайлы в папке
    all_files = os.listdir(app.config['POST_MEDIA_FOLDER'])
    video_extensions = {'.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm'}
    physical_videos = [f for f in all_files if os.path.splitext(f)[1].lower() in video_extensions]

    conn.close()

    return f"""
    <h2>Диагностика видео</h2>

    <h3>Видео в БД: {len(videos)}</h3>
    <table border="1" cellpadding="5">
        <tr>
            <th>ID</th>
            <th>Пост</th>
            <th>Автор</th>
            <th>Файл</th>
            <th>Статус</th>
            <th>Размер</th>
        </tr>
        {"".join([f"""
        <tr>
            <td>{v['id']}</td>
            <td>{v['post_id']}</td>
            <td>{v['username']}</td>
            <td>{v['filename']}</td>
            <td style="color:green">✅ Существует</td>
            <td>{v['size']} байт</td>
        </tr>
        """ for v in video_files])}
        {"".join([f"""
        <tr>
            <td>{m['id']}</td>
            <td>{m['post_id']}</td>
            <td>-</td>
            <td>{m['filename']}</td>
            <td style="color:red">❌ Файл отсутствует</td>
            <td>-</td>
        </tr>
        """ for m in missing_files])}
    </table>

    <h3>Физические видеофайлы в папке: {len(physical_videos)}</h3>
    <ul>
    {"".join([f"<li>{f}</li>" for f in physical_videos[:20]])}
    </ul>

    <h3>MIME-типы видео</h3>
    <p>Убедитесь, что сервер правильно отдает видеофайлы:</p>
    <ul>
        <li>.mp4 - video/mp4</li>
        <li>.avi - video/x-msvideo</li>
        <li>.mov - video/quicktime</li>
        <li>.wmv - video/x-ms-wmv</li>
    </ul>

    <p><a href="/home">На главную</a></p>
    """


def get_video_mime_type(filename):
    """Определяет MIME-тип видео по расширению"""
    ext = filename.split('.')[-1].lower()
    mime_types = {
        'mp4': 'video/mp4',
        'avi': 'video/x-msvideo',
        'mov': 'video/quicktime',
        'wmv': 'video/x-ms-wmv',
        'flv': 'video/x-flv',
        'webm': 'video/webm',
        'mkv': 'video/x-matroska',
        'm4v': 'video/x-m4v',
        'mpg': 'video/mpeg',
        'mpeg': 'video/mpeg'
    }
    return mime_types.get(ext, 'video/mp4')


@app.route('/fix_video_types')
def fix_video_types():
    if 'user_id' not in session:
        return redirect('/login')

    conn = get_db_connection()

    # Получаем все медиафайлы с неправильным или отсутствующим типом
    media_files = conn.execute('''
        SELECT id, filename, file_type, post_id
        FROM post_media
        WHERE file_type IS NULL OR file_type = '' OR file_type = 'unknown'
    ''').fetchall()

    fixed_count = 0
    for media in media_files:
        ext = media['filename'].split('.')[-1].lower()
        if ext in ALLOWED_IMAGE_EXTENSIONS:
            new_type = 'image'
        elif ext in ALLOWED_VIDEO_EXTENSIONS:
            new_type = 'video'
        else:
            new_type = 'unknown'

        if new_type != 'unknown':
            conn.execute('''
                UPDATE post_media 
                SET file_type = ? 
                WHERE id = ?
            ''', (new_type, media['id']))
            fixed_count += 1
            print(f"Исправлен файл {media['filename']}: {media['file_type']} -> {new_type}")

    conn.commit()
    conn.close()

    return f"<h2>Исправлено {fixed_count} записей</h2><p><a href='/debug_video'>Проверить</a></p>"


@app.route('/check_table_structure')
def check_table_structure():
    conn = get_db_connection()
    cursor = conn.cursor()

    # Получаем информацию о таблице post_media
    cursor.execute("PRAGMA table_info(post_media)")
    columns = cursor.fetchall()

    result = "Структура таблицы post_media:<br>"
    for col in columns:
        result += f"ID: {col[0]}, Имя: {col[1]}, Тип: {col[2]}, NotNull: {col[3]}, Default: {col[4]}, PK: {col[5]}<br>"

    conn.close()
    return result


@app.route('/test_video/<int:post_id>')
def test_video(post_id):
    """Прямая проверка видео для конкретного поста"""
    conn = get_db_connection()

    # Получаем медиафайлы поста
    media = conn.execute('''
        SELECT * FROM post_media WHERE post_id = ?
    ''', (post_id,)).fetchall()

    # Получаем сам пост
    post = conn.execute('''
        SELECT p.*, u.username 
        FROM posts p
        JOIN users u ON p.user_id = u.id
        WHERE p.id = ?
    ''', (post_id,)).fetchone()

    conn.close()

    html = f"<h2>Пост #{post_id}</h2>"
    html += f"<p>Автор: {post['username']}</p>"
    html += f"<p>Текст: {post['content']}</p>"
    html += f"<p>Медиафайлов в БД: {len(media)}</p>"

    for m in media:
        m_dict = dict(m)
        file_path = os.path.join(app.config['POST_MEDIA_FOLDER'], m_dict['filename'])
        file_exists = os.path.exists(file_path)

        html += f"<hr>"
        html += f"<p>ID: {m_dict['id']}</p>"
        html += f"<p>Файл: {m_dict['filename']}</p>"
        html += f"<p>Тип: {m_dict['file_type']}</p>"
        html += f"<p>Физически существует: {'✅' if file_exists else '❌'}</p>"

        if file_exists and m_dict['file_type'] == 'video':
            html += f"""
            <video width="400" controls>
                <source src="/static/uploads/posts/{m_dict['filename']}" type="video/mp4">
                Ваш браузер не поддерживает видео.
            </video>
            """

    html += f'<p><a href="/home">На главную</a></p>'
    return html


# ==================== ПРОСМОТРЩИК ДОКУМЕНТОВ ====================

def docx_to_html(filepath):
    """Конвертирует .docx в HTML"""
    doc = docx.Document(filepath)
    html_parts = []
    for para in doc.paragraphs:
        if not para.text.strip():
            html_parts.append('<br>')
            continue
        style = para.style.name.lower()
        runs_html = ''
        for run in para.runs:
            text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            css = []
            if run.bold:
                css.append('font-weight:bold')
            if run.italic:
                css.append('font-style:italic')
            if run.underline:
                css.append('text-decoration:underline')
            if run.font.size:
                css.append(f'font-size:{run.font.size.pt:.0f}pt')
            if css:
                runs_html += f'<span style="{";".join(css)}">{text}</span>'
            else:
                runs_html += text

        if 'heading 1' in style:
            html_parts.append(f'<h1>{runs_html}</h1>')
        elif 'heading 2' in style:
            html_parts.append(f'<h2>{runs_html}</h2>')
        elif 'heading 3' in style:
            html_parts.append(f'<h3>{runs_html}</h3>')
        else:
            align = para.alignment
            align_style = ''
            if align and align.name == 'CENTER':
                align_style = 'text-align:center;'
            elif align and align.name == 'RIGHT':
                align_style = 'text-align:right;'
            html_parts.append(f'<p style="margin:6px 0;{align_style}">{runs_html}</p>')

    # Таблицы
    for table in doc.tables:
        t = '<table style="border-collapse:collapse;width:100%;margin:12px 0;">'
        for row in table.rows:
            t += '<tr>'
            for cell in row.cells:
                t += f'<td style="border:1px solid #ccc;padding:6px 10px;">{cell.text.replace(chr(10), "<br>")}</td>'
            t += '</tr>'
        t += '</table>'
        html_parts.append(t)

    return '\n'.join(html_parts)


def xlsx_to_html(filepath):
    """Конвертирует .xlsx в HTML таблицу"""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    html_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        html_parts.append(f'<h3 style="margin:16px 0 8px;color:#555;">Лист: {sheet_name}</h3>')
        html_parts.append('<div style="overflow-x:auto;">')
        html_parts.append('<table style="border-collapse:collapse;min-width:100%;font-size:13px;">')
        for row in ws.iter_rows():
            html_parts.append('<tr>')
            for cell in row:
                val = '' if cell.value is None else str(cell.value)
                val = val.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                bold = 'font-weight:bold;' if cell.font and cell.font.bold else ''
                bg = ''
                if cell.fill and cell.fill.fgColor and cell.fill.fgColor.type == 'rgb':
                    color = cell.fill.fgColor.rgb
                    if color and color != '00000000' and color != 'FFFFFFFF':
                        bg = f'background:#{color[2:]};'
                html_parts.append(f'<td style="border:1px solid #ddd;padding:4px 8px;{bold}{bg}">{val}</td>')
            html_parts.append('</tr>')
        html_parts.append('</table></div>')
    return '\n'.join(html_parts)


def pptx_to_html(filepath):
    """Конвертирует .pptx в HTML слайды"""
    prs = Presentation(filepath)
    html_parts = []
    for i, slide in enumerate(prs.slides, 1):
        html_parts.append(f'''
        <div style="background:white;border:1px solid #ddd;border-radius:8px;
                    padding:30px;margin-bottom:20px;min-height:200px;position:relative;
                    box-shadow:0 2px 8px rgba(0,0,0,0.08);">
            <div style="position:absolute;top:10px;right:14px;font-size:12px;color:#aaa;">Слайд {i}</div>
        ''')
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                text_esc = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                # Определяем размер шрифта первого run
                font_size = 14
                if para.runs:
                    try:
                        sz = para.runs[0].font.size
                        if sz:
                            font_size = int(sz.pt)
                    except:
                        pass
                if font_size >= 24:
                    html_parts.append(f'<h2 style="margin:8px 0;font-size:{font_size}px;">{text_esc}</h2>')
                else:
                    html_parts.append(f'<p style="margin:4px 0;font-size:{font_size}px;">{text_esc}</p>')
        html_parts.append('</div>')
    return '\n'.join(html_parts)


def txt_to_html(filepath):
    """Конвертирует .txt в HTML"""
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()
    content = content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return f'<pre style="white-space:pre-wrap;font-family:inherit;line-height:1.6;">{content}</pre>'


@app.route('/post_file/<path:filename>')
def serve_post_file(filename):
    """Отдаёт файл поста с правильными заголовками для просмотра в браузере (не скачивания)"""
    from flask import send_from_directory
    MIME_TYPES = {
        'pdf': 'application/pdf',
        'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
        'gif': 'image/gif', 'webp': 'image/webp',
        'mp4': 'video/mp4', 'webm': 'video/webm', 'mov': 'video/quicktime',
        'txt': 'text/plain; charset=utf-8',
    }
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    mimetype = MIME_TYPES.get(ext, 'application/octet-stream')
    response = send_from_directory(app.config['POST_MEDIA_FOLDER'], filename, mimetype=mimetype)
    # Убираем заголовок скачивания — браузер должен показать файл, а не скачать
    response.headers.pop('Content-Disposition', None)
    return response


@app.route('/view_document/<int:media_id>')
def view_document(media_id):
    """Страница просмотра документа прямо на сайте"""
    if 'user_id' not in session:
        return redirect('/login')

    conn = get_db_connection()
    media = conn.execute(
        'SELECT * FROM post_media WHERE id = ?', (media_id,)
    ).fetchone()
    conn.close()

    if not media:
        return 'Файл не найден', 404

    media = dict(media)
    filepath = os.path.join(app.config['POST_MEDIA_FOLDER'], media['filename'])

    if not os.path.exists(filepath):
        return 'Файл не найден на сервере', 404

    ext = media['filename'].rsplit('.', 1)[-1].lower()
    original_name = media.get('original_filename') or media['filename']

    # PDF — отдаём напрямую в iframe
    if ext == 'pdf':
        return redirect(url_for('static', filename='uploads/posts/' + media['filename']))

    content_html = ''
    error = None

    try:
        if ext == 'doc':
            # Старый формат .doc не поддерживается python-docx — предлагаем скачать
            error = 'doc_download'
        elif ext == 'docx':
            if not DOCX_AVAILABLE:
                error = 'Библиотека python-docx не установлена.'
            else:
                content_html = docx_to_html(filepath)
        elif ext in ('xls', 'xlsx'):
            if not XLSX_AVAILABLE:
                error = 'Библиотека openpyxl не установлена.'
            else:
                content_html = xlsx_to_html(filepath)
        elif ext in ('ppt', 'pptx'):
            if not PPTX_AVAILABLE:
                error = 'Библиотека python-pptx не установлена.'
            else:
                content_html = pptx_to_html(filepath)
        elif ext == 'txt':
            content_html = txt_to_html(filepath)
        else:
            error = 'Просмотр этого типа файлов не поддерживается.'
    except Exception as e:
        error = f'Не удалось открыть файл: {str(e)}'

    if error == 'doc_download':
        content_block = f'''<div style="text-align:center;padding:40px 20px;">
            <div style="font-size:64px;margin-bottom:16px;">📄</div>
            <h2 style="margin:0 0 8px;font-size:20px;">Файл в формате .doc</h2>
            <p style="color:#666;margin-bottom:24px;font-size:14px;">
                Формат <strong>.doc</strong> (Word 97–2003) не поддерживается для просмотра в браузере.<br>
                Скачайте файл и откройте в Microsoft Word или LibreOffice.
            </p>
            <a href="/static/uploads/posts/{media['filename']}" download="{original_name}"
               style="display:inline-block;padding:12px 28px;background:#2c3e50;color:#fff;border-radius:8px;text-decoration:none;font-weight:600;font-size:15px;">
                ⬇ Скачать {original_name}
            </a>
        </div>'''
    elif error:
        content_block = f"<div class='error'>{error}</div>"
    else:
        content_block = content_html

    page_html = f'''<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{original_name}</title>
    <style>
        * {{ box-sizing: border-box; margin: 0; padding: 0; }}
        body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
               background: #f5f5f5; color: #222; }}
        .toolbar {{
            position: sticky; top: 0; z-index: 100;
            background: #2c3e50; color: white;
            padding: 12px 20px;
            display: flex; align-items: center; gap: 16px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.3);
        }}
        .toolbar .filename {{
            font-size: 15px; font-weight: 600; flex: 1;
            white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
        }}
        .toolbar a {{
            color: #7fc8f8; text-decoration: none; font-size: 13px;
            padding: 5px 12px; border: 1px solid #7fc8f8; border-radius: 4px;
            white-space: nowrap;
        }}
        .toolbar a:hover {{ background: #7fc8f8; color: #2c3e50; }}
        .toolbar .back {{ color: #ccc; font-size: 13px; cursor:pointer;
                          padding: 5px 12px; border: 1px solid #ccc; border-radius: 4px;
                          text-decoration:none; }}
        .toolbar .back:hover {{ background: #ccc; color: #2c3e50; }}
        .content {{
            max-width: 860px; margin: 30px auto; padding: 40px;
            background: white; border-radius: 8px;
            box-shadow: 0 2px 12px rgba(0,0,0,0.08);
            min-height: 400px; line-height: 1.7;
        }}
        .error {{
            color: #c0392b; background: #fdecea; padding: 20px;
            border-radius: 6px; border: 1px solid #f5c6cb;
        }}
        h1,h2,h3 {{ margin: 16px 0 8px; color: #1a1a1a; }}
        table {{ margin: 12px 0; }}
    </style>
</head>
<body>
    <div class="toolbar">
        <span class="filename">📄 {original_name}</span>
        <a href="/static/uploads/posts/{media['filename']}" download="{original_name}">⬇ Скачать</a>
    </div>
    <div class="content">
        {content_block}
    </div>
</body>
</html>'''

    return page_html


@app.route('/view_document_content/<int:media_id>')
def view_document_content(media_id):
    """Возвращает только HTML-содержимое документа для показа в модальном окне"""
    if 'user_id' not in session:
        return jsonify({'error': 'Не авторизован'}), 401

    conn = get_db_connection()
    media = conn.execute('SELECT * FROM post_media WHERE id = ?', (media_id,)).fetchone()
    conn.close()

    if not media:
        return '<p style="color:red">Файл не найден</p>', 404

    media = dict(media)
    filepath = os.path.join(app.config['POST_MEDIA_FOLDER'], media['filename'])

    if not os.path.exists(filepath):
        return '<p style="color:red">Файл не найден на сервере</p>', 404

    ext = media['filename'].rsplit('.', 1)[-1].lower()
    original_name = media.get('original_filename') or media['filename']

    try:
        if ext == 'pdf':
            # Отдаём PDF через отдельный маршрут с правильными заголовками
            file_url = url_for('serve_post_file', filename=media['filename'])
            return f'''<iframe src="{file_url}" style="width:100%;height:100%;border:none;min-height:80vh;display:block;"></iframe>'''
        elif ext in ('doc', 'docx') and DOCX_AVAILABLE:
            html = docx_to_html(filepath)
        elif ext in ('xls', 'xlsx') and XLSX_AVAILABLE:
            html = xlsx_to_html(filepath)
        elif ext in ('ppt', 'pptx') and PPTX_AVAILABLE:
            html = pptx_to_html(filepath)
        elif ext == 'txt':
            html = txt_to_html(filepath)
        else:
            file_url = url_for('serve_post_file', filename=media['filename'])
            return f'''<div style="text-align:center;padding:40px;">
                <div style="font-size:64px;">📎</div>
                <p style="margin:16px 0;font-size:16px;">{original_name}</p>
                <a href="{file_url}" download="{original_name}"
                   style="display:inline-block;padding:10px 24px;background:#007bff;color:white;
                          border-radius:6px;text-decoration:none;">⬇ Скачать файл</a>
            </div>'''
        return f'''<div style="padding:24px;font-family:-apple-system,BlinkMacSystemFont,Segoe UI,sans-serif;
                           line-height:1.7;max-width:860px;margin:0 auto;">{html}</div>'''
    except Exception as e:
        return f'''<div style="padding:24px;color:#c0392b;background:#fdecea;border-radius:6px;margin:20px;">
            Не удалось открыть файл: {str(e)}</div>''', 500


# ==================== МЕССЕНДЖЕР ====================

def ensure_messenger_tables():
    conn = get_db_connection()
    conn.execute('''
        CREATE TABLE IF NOT EXISTS conversations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user1_id INTEGER NOT NULL,
            user2_id INTEGER NOT NULL,
            last_message_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user1_id) REFERENCES users(id),
            FOREIGN KEY (user2_id) REFERENCES users(id),
            UNIQUE(user1_id, user2_id)
        )
    ''')
    conn.execute('''
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            conversation_id INTEGER NOT NULL,
            sender_id INTEGER NOT NULL,
            text TEXT NOT NULL,
            is_read INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (conversation_id) REFERENCES conversations(id) ON DELETE CASCADE,
            FOREIGN KEY (sender_id) REFERENCES users(id)
        )
    ''')
    conn.commit()
    conn.close()


def get_or_create_conversation(user1_id, user2_id):
    conn = get_db_connection()
    a, b = sorted([user1_id, user2_id])
    conv = conn.execute(
        'SELECT id FROM conversations WHERE user1_id = ? AND user2_id = ?', (a, b)
    ).fetchone()
    if conv:
        conv_id = conv['id']
    else:
        cursor = conn.cursor()
        cursor.execute(
            'INSERT INTO conversations (user1_id, user2_id) VALUES (?, ?)', (a, b)
        )
        conv_id = cursor.lastrowid
        conn.commit()
    conn.close()
    return conv_id


def get_conversations_list(user_id):
    conn = get_db_connection()
    rows = rows_to_dicts(conn.execute('''
        SELECT
            c.id as conv_id,
            CASE WHEN c.user1_id = ? THEN c.user2_id ELSE c.user1_id END as partner_id,
            u.username as partner_username,
            COALESCE(up.full_name, u.username) as partner_name,
            COALESCE(up.avatar, 'default_avatar.png') as partner_avatar,
            (SELECT text FROM messages WHERE conversation_id = c.id ORDER BY created_at DESC LIMIT 1) as last_message,
            (SELECT COUNT(*) FROM messages WHERE conversation_id = c.id AND sender_id != ? AND is_read = 0) as unread_count
        FROM conversations c
        JOIN users u ON u.id = CASE WHEN c.user1_id = ? THEN c.user2_id ELSE c.user1_id END
        LEFT JOIN user_profiles up ON up.user_id = u.id
        WHERE c.user1_id = ? OR c.user2_id = ?
        ORDER BY c.last_message_at DESC
    ''', (user_id, user_id, user_id, user_id, user_id)).fetchall())
    conn.close()
    return rows


@app.route('/messenger')
def messenger():
    if 'user_id' not in session:
        return redirect('/login')
    ensure_messenger_tables()
    user_id = session['user_id']
    conversations = get_conversations_list(user_id)
    total_unread = sum(c['unread_count'] for c in conversations)
    return render_template('messenger.html',
                           conversations=conversations,
                           total_unread=total_unread,
                           active_conv=None,
                           messages=[],
                           partner=None,
                           partner_id=None,
                           user_id=user_id)


@app.route('/messenger/<int:partner_id>')
def messenger_chat(partner_id):
    if 'user_id' not in session:
        return redirect('/login')
    ensure_messenger_tables()
    user_id = session['user_id']
    if user_id == partner_id:
        flash('Нельзя писать самому себе', 'error')
        return redirect('/messenger')
    conv_id = get_or_create_conversation(user_id, partner_id)
    conn = get_db_connection()
    conn.execute('''
        UPDATE messages SET is_read = 1
        WHERE conversation_id = ? AND sender_id != ?
    ''', (conv_id, user_id))
    conn.commit()
    chat_messages = rows_to_dicts(conn.execute('''
        SELECT m.id, m.sender_id, m.text, m.created_at, u.username as sender_username
        FROM messages m
        JOIN users u ON m.sender_id = u.id
        WHERE m.conversation_id = ?
        ORDER BY m.created_at ASC
        LIMIT 200
    ''', (conv_id,)).fetchall())
    partner = row_to_dict(conn.execute('''
        SELECT u.id, u.username,
               COALESCE(up.full_name, u.username) as full_name,
               COALESCE(up.avatar, 'default_avatar.png') as avatar
        FROM users u
        LEFT JOIN user_profiles up ON up.user_id = u.id
        WHERE u.id = ?
    ''', (partner_id,)).fetchone())
    is_blocked_by_me = conn.execute(
        'SELECT id FROM blacklist WHERE blocker_id = ? AND blocked_id = ?',
        (user_id, partner_id)
    ).fetchone() is not None
    is_blocked_by_them = conn.execute(
        'SELECT id FROM blacklist WHERE blocker_id = ? AND blocked_id = ?',
        (partner_id, user_id)
    ).fetchone() is not None
    conn.close()
    conversations = get_conversations_list(user_id)
    total_unread = sum(c['unread_count'] for c in conversations)
    return render_template('messenger.html',
                           conversations=conversations,
                           total_unread=total_unread,
                           active_conv=conv_id,
                           messages=chat_messages,
                           partner=partner,
                           partner_id=partner_id,
                           current_user_id=user_id,
                           user_id=user_id,
                           is_blocked_by_me=is_blocked_by_me,
                           is_blocked_by_them=is_blocked_by_them)


@app.route('/messenger/send_ajax', methods=['POST'])
def messenger_send_ajax():
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    ensure_messenger_tables()
    user_id = session['user_id']
    data = request.get_json()
    partner_id = data.get('partner_id')
    text = (data.get('text') or '').strip()
    if not text or not partner_id:
        return jsonify({'success': False, 'error': 'Пустое сообщение'})
    conv_id = get_or_create_conversation(user_id, int(partner_id))
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute(
        'INSERT INTO messages (conversation_id, sender_id, text) VALUES (?, ?, ?)',
        (conv_id, user_id, text)
    )
    msg_id = cursor.lastrowid
    conn.execute(
        'UPDATE conversations SET last_message_at = CURRENT_TIMESTAMP WHERE id = ?',
        (conv_id,)
    )
    conn.commit()
    msg = row_to_dict(conn.execute('SELECT * FROM messages WHERE id = ?', (msg_id,)).fetchone())
    conn.close()
    return jsonify({'success': True, 'message': msg})


@app.route('/messenger/poll/<int:conv_id>')
def messenger_poll(conv_id):
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    ensure_messenger_tables()
    user_id = session['user_id']
    last_id = request.args.get('last_id', 0, type=int)
    conn = get_db_connection()
    conv = conn.execute(
        'SELECT * FROM conversations WHERE id = ? AND (user1_id = ? OR user2_id = ?)',
        (conv_id, user_id, user_id)
    ).fetchone()
    if not conv:
        conn.close()
        return jsonify({'success': False, 'error': 'Нет доступа'})
    new_messages = rows_to_dicts(conn.execute('''
        SELECT m.id, m.sender_id, m.text, m.created_at, u.username as sender_username
        FROM messages m
        JOIN users u ON m.sender_id = u.id
        WHERE m.conversation_id = ? AND m.id > ?
        ORDER BY m.created_at ASC
    ''', (conv_id, last_id)).fetchall())
    if new_messages:
        conn.execute('''
            UPDATE messages SET is_read = 1
            WHERE conversation_id = ? AND sender_id != ? AND id > ?
        ''', (conv_id, user_id, last_id))
        conn.commit()
    total_unread = conn.execute('''
        SELECT COUNT(*) as count FROM messages m
        JOIN conversations c ON m.conversation_id = c.id
        WHERE (c.user1_id = ? OR c.user2_id = ?)
          AND m.sender_id != ? AND m.is_read = 0
          AND m.conversation_id != ?
    ''', (user_id, user_id, user_id, conv_id)).fetchone()['count']
    conn.close()
    return jsonify({'success': True, 'messages': new_messages, 'total_unread': total_unread})


@app.route('/messenger/unread_count')
def messenger_unread_count():
    if 'user_id' not in session:
        return jsonify({'count': 0})
    ensure_messenger_tables()
    user_id = session['user_id']
    conn = get_db_connection()
    count = conn.execute('''
        SELECT COUNT(*) as count FROM messages m
        JOIN conversations c ON m.conversation_id = c.id
        WHERE (c.user1_id = ? OR c.user2_id = ?) AND m.sender_id != ? AND m.is_read = 0
    ''', (user_id, user_id, user_id)).fetchone()['count']
    conn.close()
    return jsonify({'count': count})


@app.route('/messenger/search_users')
def messenger_search_users():
    if 'user_id' not in session:
        return jsonify({'users': []})
    ensure_messenger_tables()
    user_id = session['user_id']
    q = request.args.get('q', '').strip()
    if len(q) < 2:
        return jsonify({'users': []})
    conn = get_db_connection()
    users = rows_to_dicts(conn.execute('''
        SELECT u.id, u.username,
               COALESCE(up.full_name, u.username) as full_name,
               COALESCE(up.avatar, 'default_avatar.png') as avatar
        FROM users u
        LEFT JOIN user_profiles up ON up.user_id = u.id
        WHERE (u.username LIKE ? OR up.full_name LIKE ?)
          AND u.id != ? AND u.is_banned = 0
        LIMIT 10
    ''', (f'%{q}%', f'%{q}%', user_id)).fetchall())
    conn.close()
    return jsonify({'users': users})


@app.route('/messenger/delete_conv/<int:partner_id>', methods=['POST'])
def messenger_delete_conv(partner_id):
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    ensure_messenger_tables()
    user_id = session['user_id']
    conn = get_db_connection()
    conv = conn.execute(
        'SELECT id FROM conversations WHERE (user1_id=? AND user2_id=?) OR (user1_id=? AND user2_id=?)',
        (user_id, partner_id, partner_id, user_id)
    ).fetchone()
    if conv:
        conn.execute('DELETE FROM messages WHERE conversation_id = ?', (conv['id'],))
        conn.execute('DELETE FROM conversations WHERE id = ?', (conv['id'],))
        conn.commit()
    conn.close()
    return jsonify({'success': True})


@app.route('/messenger/clear_history/<int:partner_id>', methods=['POST'])
def messenger_clear_history(partner_id):
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    ensure_messenger_tables()
    user_id = session['user_id']
    conn = get_db_connection()
    conv = conn.execute(
        'SELECT id FROM conversations WHERE (user1_id=? AND user2_id=?) OR (user1_id=? AND user2_id=?)',
        (user_id, partner_id, partner_id, user_id)
    ).fetchone()
    if not conv:
        conn.close()
        return jsonify({'success': False, 'error': 'Диалог не найден'})
    conn.execute('DELETE FROM messages WHERE conversation_id = ?', (conv['id'],))
    conn.commit()
    conn.close()
    return jsonify({'success': True})


@app.route('/messenger/block/<int:partner_id>', methods=['POST'])
def messenger_block(partner_id):
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    user_id = session['user_id']
    if user_id == partner_id:
        return jsonify({'success': False, 'error': 'Нельзя заблокировать себя'})
    conn = get_db_connection()
    existing = conn.execute(
        'SELECT id FROM blacklist WHERE blocker_id = ? AND blocked_id = ?',
        (user_id, partner_id)
    ).fetchone()
    if not existing:
        conn.execute(
            'INSERT INTO blacklist (blocker_id, blocked_id, reason) VALUES (?, ?, ?)',
            (user_id, partner_id, 'Заблокирован из мессенджера')
        )
        conn.commit()
    conn.close()
    return jsonify({'success': True})


@app.route('/messenger/unblock/<int:partner_id>', methods=['POST'])
def messenger_unblock(partner_id):
    if 'user_id' not in session:
        return jsonify({'success': False}), 401
    user_id = session['user_id']
    conn = get_db_connection()
    conn.execute(
        'DELETE FROM blacklist WHERE blocker_id = ? AND blocked_id = ?',
        (user_id, partner_id)
    )
    conn.commit()
    conn.close()
    return jsonify({'success': True})


if __name__ == '__main__':
    create_tables()

    # Запускаем Telegram-бота в фоновом потоке
    import os

    if not app.debug or os.environ.get('WERKZEUG_RUN_MAIN') == 'true':
        bot_thread = threading.Thread(target=run_telegram_bot, daemon=True)
        bot_thread.start()

    app.run(debug=True, host='0.0.0.0', port=5555, use_reloader=False)